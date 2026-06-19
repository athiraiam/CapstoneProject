"""
update_pptx_agents.py
Rebuilds the full presentation deck including all agent/skill/hooks slides.
Run: python update_pptx_agents.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import datetime

DARK_BLUE  = RGBColor(0,   51,  102)
MID_BLUE   = RGBColor(0,   102, 204)
LIGHT_BLUE = RGBColor(204, 229, 255)
WHITE      = RGBColor(255, 255, 255)
GOLD       = RGBColor(255, 193,   7)
DARK_GRAY  = RGBColor(50,   50,  50)
GREEN      = RGBColor(0,   140,  64)
PURPLE     = RGBColor(100,  0,  160)
ORANGE     = RGBColor(200,  90,   0)
RED        = RGBColor(180,   0,   0)
TEAL       = RGBColor(0,   120, 120)
TODAY = datetime.date.today().strftime("%B %d, %Y")

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, text, left, top, w, h, size=14, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text          = text
    r.font.size     = Pt(size)
    r.font.bold     = bold
    r.font.italic   = italic
    r.font.color.rgb = color
    return box


def rect(slide, left, top, w, h, fill, line=False):
    s = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if not line:
        s.line.fill.background()
    return s


def hdr(slide, title, color=DARK_BLUE):
    rect(slide, 0, 0, 13.33, 1.2, color)
    tb(slide, title, 0.4, 0.18, 12.5, 0.9, size=26, bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BLUE)
rect(s, 0, 3.0, 13.33, 0.06, GOLD)
tb(s, "🏦  Internal Bank Employee Assistant", 0.5, 0.7, 12.3, 1.4,
   size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Multi-Agent AI System — Capstone Project", 0.5, 2.1, 12.3, 0.7,
   size=19, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
tb(s, f"NexaBank Internal Division     |     {TODAY}", 0.5, 3.2, 12.3, 0.6,
   size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
tb(s, "Python  •  Streamlit  •  7 Agents  •  3 Skills  •  6 Hooks",
   0.5, 4.0, 12.3, 0.6, size=13, color=GOLD, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — Business Problem
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
hdr(s, "Business Problem")
problems = [
    "Bank employees must reference 5+ policies daily — manual search takes 5-10 minutes",
    "Inconsistent answers lead to compliance risks and customer service failures",
    "Policy documents are 50+ lines each — finding the right section is error-prone",
    "No single system to classify, search, summarize, escalate, and audit policy queries",
    "New employees struggle to navigate complex policy documents independently",
]
for i, p in enumerate(problems):
    rect(s, 0.4, 1.4 + i*0.9, 12.3, 0.75, LIGHT_BLUE)
    tb(s, f"⚠  {p}", 0.55, 1.52 + i*0.9, 12.0, 0.55, size=13, color=DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution Overview
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
hdr(s, "Solution Overview", MID_BLUE)
tb(s, "A fully orchestrated multi-agent AI system that classifies, searches, summarizes, escalates, and audits — all from an offline knowledge base.",
   0.4, 1.25, 12.5, 0.6, size=13, color=DARK_BLUE, italic=True)

features = [
    ("🤖 OrchestratorAgent", "Coordinates all 7 subagents"),
    ("🔍 Intent Classifier",  "Routes to the right agent"),
    ("📚 Policy Search",      "Keyword scoring on 5 policies"),
    ("📝 Multi-Policy",       "Merges answers from 2 policies"),
    ("📋 Summary Agent",      "5-bullet policy summaries"),
    ("📧 Escalation Agent",   "Auto-drafts escalation emails"),
    ("🛡️ Audit Logger",       "Logs every query silently"),
    ("⚠️ Self-Healing",       "FallbackAgent auto-recovers"),
    ("🔗 6 Hooks",            "Pre/post query lifecycle"),
]
for i, (title, desc) in enumerate(features):
    left = 0.3 + (i % 3) * 4.3
    top  = 2.0 + (i // 3) * 1.7
    rect(s, left, top, 4.1, 1.5, DARK_BLUE)
    tb(s, title, left+0.1, top+0.1,  3.9, 0.55, size=12, bold=True, color=GOLD)
    tb(s, desc,  left+0.1, top+0.65, 3.9, 0.7,  size=11, color=WHITE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — Orchestrator Pipeline
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
hdr(s, "OrchestratorAgent — Full Pipeline")

steps = [
    ("1", "pre_query_hook",         "Validate & sanitize input",        MID_BLUE),
    ("2", "QueryClassifierAgent",   "Detect intent",                    GREEN),
    ("3", "Route Decision",         "Orchestrator selects subagent",    DARK_BLUE),
    ("4", "SubAgent Executes",      "Answer produced",                  PURPLE),
    ("5", "AuditLoggerAgent",       "Log silently in background",       TEAL),
    ("6", "compliance_alert_hook",  "Flag risk keywords if any",        RED),
    ("7", "escalation_hook",        "Log unresolved queries",           ORANGE),
    ("8", "post_answer_hook",       "Log response time + outcome",      MID_BLUE),
]
for i, (num, name, detail, color) in enumerate(steps):
    top = 1.35 + i * 0.72
    rect(s, 0.3, top, 0.55, 0.58, color)
    tb(s, num,    0.3,  top+0.1, 0.55, 0.42, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, name,   1.0,  top,     4.0,  0.35, size=12, bold=True, color=DARK_BLUE)
    tb(s, detail, 1.0,  top+0.32,4.0,  0.3,  size=10, color=DARK_GRAY, italic=True)
    # arrow
    if i < 7:
        tb(s, "↓", 0.42, top+0.58, 0.4, 0.2, size=10, color=color, align=PP_ALIGN.CENTER)

tb(s, "Result returned to Streamlit UI with: answer + source badge + agent chip + compliance alert (if any)",
   5.5, 1.35, 7.5, 5.7, size=12, color=DARK_BLUE)

# right-side box
rect(s, 5.5, 1.35, 7.4, 5.7, LIGHT_BLUE)
tb(s, "Result Dictionary Fields:", 5.65, 1.5, 7.1, 0.5, size=13, bold=True, color=DARK_BLUE)
fields = ["answer", "policy_name / policy_names", "found (bool)",
          "intent", "agent_used", "escalated", "email_draft",
          "compliance_alert", "session_id"]
for i, f in enumerate(fields):
    tb(s, f"  ·  {f}", 5.65, 2.05 + i*0.52, 7.1, 0.45, size=11, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — 7 Subagents
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
hdr(s, "7 Subagents — Roles & Responsibilities")

agents = [
    ("QueryClassifierAgent",  "Detects intent before search",      "policy_query / summary / escalation",    GREEN),
    ("PolicySearchAgent",     "Single-policy keyword search",       "Best match from 5 policies",            MID_BLUE),
    ("MultiPolicyAgent",      "Spans 2 policy domains",            "Merges top-2 policy snippets",          PURPLE),
    ("SummaryAgent",          "On-demand policy overview",         "5-bullet summary on request",           ORANGE),
    ("EscalationAgent",       "Routes to correct team",            "Drafts escalation email automatically", RED),
    ("FallbackAgent",         "Self-healing on no match",          "Partial recovery + escalation trigger", RGBColor(120,60,0)),
    ("AuditLoggerAgent",      "Background audit logging",          "Logs every query silently to file",     TEAL),
]
for i, (name, role, detail, color) in enumerate(agents):
    left = 0.3 + (i % 4) * 3.22
    top  = 1.4 + (i // 4) * 2.5
    rect(s, left, top, 3.1, 2.2, color)
    tb(s, name,   left+0.1, top+0.1,  2.9, 0.6, size=11, bold=True, color=GOLD)
    tb(s, role,   left+0.1, top+0.7,  2.9, 0.6, size=10, color=WHITE)
    tb(s, detail, left+0.1, top+1.3,  2.9, 0.75, size=9, color=LIGHT_BLUE, italic=True)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — 3 Skills
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
hdr(s, "Skills — Reusable Functions Across Agents")

skills = [
    (
        "search_skill.py", GREEN,
        [
            ("search_policy()",     "Keyword scores all 5 policies, returns best match + snippet"),
            ("extract_keywords()",  "Pulls meaningful words from question for classification"),
            ("cross_reference()",   "Scores ALL policies → enables multi-policy detection"),
        ],
        ["PolicySearchAgent", "MultiPolicyAgent", "SummaryAgent", "FallbackAgent"]
    ),
    (
        "format_skill.py", PURPLE,
        [
            ("format_answer()",              "Cleans raw policy text into readable markdown"),
            ("summarize_policy()",           "5-bullet summary of any policy document"),
            ("generate_escalation_email()",  "Drafts pre-filled escalation email body"),
        ],
        ["PolicySearchAgent", "SummaryAgent", "EscalationAgent"]
    ),
    (
        "audit_skill.py", TEAL,
        [
            ("log_query()",              "Appends to query_audit.log"),
            ("log_unresolved()",         "Appends to unresolved_queries.log"),
            ("log_session_event()",      "Logs session lifecycle events"),
            ("detect_compliance_risk()", "Returns True if risk keywords found"),
        ],
        ["AuditLoggerAgent", "All 6 Hooks"]
    ),
]
for i, (filename, color, functions, used_by) in enumerate(skills):
    left = 0.3 + i * 4.3
    rect(s, left, 1.35, 4.1, 5.6, color)
    tb(s, filename, left+0.1, 1.45, 3.9, 0.55, size=14, bold=True, color=GOLD)
    tb(s, "Functions:", left+0.1, 2.05, 3.9, 0.35, size=10, bold=True, color=WHITE)
    for j, (fn, desc) in enumerate(functions):
        tb(s, fn,   left+0.15, 2.45 + j*0.9, 3.8, 0.35, size=10, bold=True, color=LIGHT_BLUE)
        tb(s, desc, left+0.15, 2.78 + j*0.9, 3.8, 0.45, size=9,  color=WHITE, italic=True)
    tb(s, "Used by:", left+0.1, 5.4, 3.9, 0.3, size=10, bold=True, color=WHITE)
    tb(s, " · ".join(used_by), left+0.1, 5.7, 3.9, 0.8, size=9, color=LIGHT_BLUE, italic=True)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — 6 Hooks
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
hdr(s, "Hooks — Lifecycle Callbacks")
tb(s, "Hooks fire automatically at specific pipeline points. They handle cross-cutting concerns without polluting agent logic.",
   0.4, 1.2, 12.5, 0.5, size=12, color=DARK_BLUE, italic=True)

hooks = [
    ("pre_query_hook",        "BEFORE any agent",         "Validates input, rejects empty questions, logs incoming query",                       MID_BLUE),
    ("post_answer_hook",      "AFTER pipeline completes", "Logs matched policy, response time, pass/fallback status",                            MID_BLUE),
    ("session_start_hook",    "App loads (once)",         "Logs session start with unique session ID",                                          GREEN),
    ("session_end_hook",      "Employee clicks Clear",    "Logs session closure + total query count for the session",                           GREEN),
    ("escalation_hook",       "When found=False",         "Writes unmatched question to unresolved_queries.log for compliance team review",     RED),
    ("compliance_alert_hook", "Risk keyword detected",    "Logs flagged query, returns yellow warning banner text for the UI",                  ORANGE),
]
for i, (name, trigger, responsibility, color) in enumerate(hooks):
    top = 1.85 + i * 0.88
    rect(s, 0.3, top, 12.7, 0.75, LIGHT_BLUE if i % 2 == 0 else RGBColor(220, 235, 255))
    rect(s, 0.3, top, 0.18, 0.75, color)
    tb(s, name,           0.6,  top+0.08, 2.8, 0.35, size=11, bold=True,  color=DARK_BLUE)
    tb(s, f"→ {trigger}", 0.6,  top+0.42, 2.8, 0.28, size=9,  italic=True, color=DARK_GRAY)
    tb(s, responsibility, 3.5,  top+0.18, 9.3, 0.45, size=11, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — Routing Logic (Agent Decision Tree)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
hdr(s, "Agent Routing — Decision Tree")

# Draw decision flow
nodes = [
    (5.7, 1.3,  2.0, 0.65, "Employee Question",               DARK_BLUE),
    (5.7, 2.2,  2.0, 0.65, "QueryClassifier\nAgent",          GREEN),
    (0.3, 3.4,  2.5, 0.65, "summary intent",                  PURPLE),
    (2.9, 3.4,  2.5, 0.65, "escalation intent",               RED),
    (5.5, 3.4,  2.6, 0.65, "policy_query\n(multi-policy)",    MID_BLUE),
    (8.2, 3.4,  2.5, 0.65, "policy_query\n(single policy)",   MID_BLUE),
    (10.8,3.4,  2.2, 0.65, "out of scope",                    DARK_GRAY),
    (0.3, 4.6,  2.5, 0.65, "SummaryAgent",                    PURPLE),
    (2.9, 4.6,  2.5, 0.65, "EscalationAgent",                 RED),
    (5.5, 4.6,  2.6, 0.65, "MultiPolicyAgent",                MID_BLUE),
    (8.2, 4.6,  2.5, 0.65, "PolicySearchAgent",               MID_BLUE),
    (10.8,4.6,  2.2, 0.65, "Fallback\nMessage",               DARK_GRAY),
    (8.2, 5.7,  2.5, 0.65, "FallbackAgent\n(self-healing)",   ORANGE),
]
for left, top, w, h, label, color in nodes:
    rect(s, left, top, w, h, color)
    tb(s, label, left+0.05, top+0.08, w-0.1, h-0.1, size=10, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)

# Arrows (text)
tb(s, "↓", 6.55, 1.95, 0.4, 0.3, size=14, bold=True, color=DARK_BLUE)
for x in [1.35, 3.95, 6.6, 9.3, 11.7]:
    tb(s, "↓", x, 4.05, 0.4, 0.3, size=14, bold=True, color=DARK_GRAY)
tb(s, "↓ no match", 9.0, 5.28, 1.5, 0.42, size=9, italic=True, color=ORANGE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — Self-Healing Workflow
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
hdr(s, "Self-Healing Workflow — FallbackAgent")

tb(s, "When PolicySearchAgent returns found=False, the Orchestrator automatically chains FallbackAgent and EscalationAgent without any employee action.",
   0.4, 1.25, 12.5, 0.6, size=13, color=DARK_BLUE, italic=True)

steps = [
    (MID_BLUE,  "1. PolicySearchAgent",      "Keyword search returns score=0 → found=False"),
    (ORANGE,    "2. OrchestratorAgent",       "Detects found=False → automatically invokes FallbackAgent"),
    (ORANGE,    "3. FallbackAgent",           "Tries cross_reference at low threshold for partial match"),
    (RED,       "4. escalation_hook",         "Fires: writes question to unresolved_queries.log"),
    (RED,       "5. EscalationAgent (opt.)",  "If employee asks for contact — provides team + email draft"),
    (TEAL,      "6. AuditLoggerAgent",        "Logs fallback event with session ID to query_audit.log"),
    (MID_BLUE,  "7. UI displays",             "Warning box + contact details + agent chip shown to employee"),
]
for i, (color, title, detail) in enumerate(steps):
    top = 2.0 + i * 0.72
    rect(s, 0.3, top, 0.5, 0.6, color)
    tb(s, str(i+1), 0.3, top+0.1, 0.5, 0.42, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, title,    1.0, top,      4.5, 0.35, size=12, bold=True, color=DARK_BLUE)
    tb(s, detail,   1.0, top+0.35, 11.8,0.3,  size=11, color=DARK_GRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — Governance
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
hdr(s, "Governance Framework")

pillars = [
    ("🔒 No PII",            ["No customer data processed", "Queries cleared on close", "No external API calls"],         DARK_BLUE),
    ("📋 Source Control",    ["Every answer cites policy", "Verbatim from approved text", "No hallucination possible"],   MID_BLUE),
    ("🛡️ Agent Governance",  ["Each agent: read-only", "Only AuditLoggerAgent writes", "Append-only log files"],          GREEN),
    ("⚠️ Risk Detection",    ["14 compliance keywords", "compliance_alert_hook fires", "Logged + warning shown in UI"],   RED),
    ("🔄 Audit Trail",       ["query_audit.log", "unresolved_queries.log", "session_events.log"],                         TEAL),
    ("👥 Access controls",   ["Internal network only", "Policy owners manage .txt", "IT controls deployment"],            ORANGE),
]
for i, (title, points, color) in enumerate(pillars):
    left = 0.3 + (i % 3) * 4.3
    top  = 1.4 + (i // 3) * 2.7
    rect(s, left, top, 4.1, 2.4, LIGHT_BLUE)
    rect(s, left, top, 4.1, 0.55, color)
    tb(s, title, left+0.1, top+0.08, 3.9, 0.42, size=13, bold=True, color=WHITE)
    for j, pt in enumerate(points):
        tb(s, f"• {pt}", left+0.15, top+0.65+j*0.55, 3.8, 0.48, size=11, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — Test Results
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
hdr(s, "Testing & Evaluation Results")

metrics = [("47", "Total Tests"), ("47", "Passed"), ("0", "Failed"), ("100%", "Pass Rate")]
for i, (val, label) in enumerate(metrics):
    left = 0.5 + i * 3.1
    c = GREEN if i == 3 else MID_BLUE
    rect(s, left, 1.4, 2.8, 1.8, c)
    tb(s, val,   left+0.1, 1.5,  2.6, 1.0, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, label, left+0.1, 2.5,  2.6, 0.5, size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

categories = [
    ("Policy Search (5 domains)", "18/18", GREEN),
    ("QueryClassifierAgent",      "6/6",   GREEN),
    ("MultiPolicyAgent",          "3/3",   GREEN),
    ("SummaryAgent",              "3/3",   GREEN),
    ("EscalationAgent",           "4/4",   GREEN),
    ("FallbackAgent",             "5/5",   GREEN),
    ("Audit Logger & Hooks",      "5/5",   GREEN),
    ("Compliance Risk",           "3/3",   GREEN),
]
tb(s, "By Category:", 0.4, 3.5, 4.0, 0.4, size=13, bold=True, color=DARK_BLUE)
for i, (cat, score, color) in enumerate(categories):
    left = 0.4 + (i % 4) * 3.2
    top  = 3.95 + (i // 4) * 0.75
    tb(s, f"✅  {cat}: {score}", left, top, 3.1, 0.6, size=11, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — Observability
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
hdr(s, "Observability & Traceability", MID_BLUE)

items = [
    ("🔍 Agent Chip on every answer",    "UI shows which agent answered + which intent was detected"),
    ("📋 Source Badge",                  "Every answer shows the policy name and code (e.g. LP-001)"),
    ("⚠️ Compliance Alert Banner",      "Yellow warning shown in UI when risk keywords detected"),
    ("📧 Email Draft Expander",          "EscalationAgent shows draft email inline in chat"),
    ("📝 query_audit.log",               "Timestamp | Session | Status (MATCHED/FALLBACK) | Policy | Question"),
    ("📝 unresolved_queries.log",        "Every no-match query logged for compliance team review"),
    ("📝 session_events.log",            "Session start/end, all hooks, compliance alerts, response times"),
    ("🚫 No Black Box",                  "Deterministic rules — every answer traceable to a policy line"),
]
for i, (label, desc) in enumerate(items):
    top = 1.4 + i * 0.74
    c = LIGHT_BLUE if i % 2 == 0 else RGBColor(220, 235, 255)
    rect(s, 0.3, top, 12.7, 0.62, c)
    tb(s, label, 0.5,  top+0.12, 3.8, 0.42, size=12, bold=True,  color=DARK_BLUE)
    tb(s, desc,  4.4,  top+0.12, 8.5, 0.42, size=12, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — Business Impact
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
hdr(s, "Business Impact")

impacts = [
    ("⏱️ 95% Faster",         "Policy lookup: 5-10 min → under 30 sec"),
    ("✅ Consistent",          "All employees get identical approved answers"),
    ("🤖 7 Agents",            "Handles search, summary, escalation, audit"),
    ("📉 Reduced Risk",        "Compliance risk auto-flagged in real time"),
    ("📧 Auto Escalation",     "Draft email generated, no manual drafting"),
    ("🔒 Zero Data Risk",      "No PII, no internet, no external APIs"),
]
for i, (title, detail) in enumerate(impacts):
    left = 0.4 + (i % 2) * 6.4
    top  = 1.4 + (i // 2) * 1.9
    rect(s, left, top, 6.0, 1.7, LIGHT_BLUE)
    rect(s, left, top, 6.0, 0.55, DARK_BLUE)
    tb(s, title,  left+0.15, top+0.08, 5.7, 0.42, size=15, bold=True, color=WHITE)
    tb(s, detail, left+0.15, top+0.7,  5.7, 0.75, size=13, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 14 — Thank You
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BLUE)
rect(s, 0, 3.1, 13.33, 0.06, GOLD)
tb(s, "🏦", 6.2, 0.5, 1.0, 1.1, size=42, align=PP_ALIGN.CENTER, color=WHITE)
tb(s, "Thank You", 0.5, 1.6, 12.3, 1.1, size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Internal Bank Employee Assistant — Multi-Agent Capstone Project",
   0.5, 2.6, 12.3, 0.6, size=15, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
tb(s, "Questions & Discussion",
   0.5, 3.4, 12.3, 0.8, size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s, "7 Agents  •  3 Skills  •  6 Hooks  •  3 Log Files  •  47 Tests  •  100% Pass Rate",
   0.5, 4.5, 12.3, 0.6, size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

prs.save("/home/labuser/Day10/deliverables/05_Presentation_Deck.pptx")
print("✅ Presentation Deck updated with agent/skill/hooks slides.")
