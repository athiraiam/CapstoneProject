"""
make_architecture_ppt.py
Generates a dedicated Architecture Diagram presentation for the
Internal Bank Employee Assistant — Multi-Agent System.
6 slides:
  1. Cover
  2. Full System Architecture (layered overview)
  3. Agent Layer — Routing & Decision Tree
  4. End-to-End Data Flow (one query, 9 steps)
  5. Skills × Agents Usage Matrix + Hook Lifecycle
  6. Self-Healing Workflow Deep-Dive
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches as I
import copy
from lxml import etree

OUT = "/home/labuser/Day10/deliverables/07_Architecture_Diagram.pptx"

prs = Presentation()
prs.slide_width  = I(13.33)
prs.slide_height = I(7.5)
BLANK = prs.slide_layouts[6]

# ── Palette ────────────────────────────────────────────────────────────────
C = {
    "dark_navy":   RGBColor(10,  25,  70),
    "navy":        RGBColor(0,   51, 102),
    "blue":        RGBColor(0,  102, 204),
    "sky":         RGBColor(204,229, 255),
    "white":       RGBColor(255,255, 255),
    "gold":        RGBColor(255,193,   7),
    "green":       RGBColor(0,  130,  60),
    "lt_green":    RGBColor(210,240, 220),
    "purple":      RGBColor(100,  0, 160),
    "lt_purple":   RGBColor(235,215, 255),
    "red":         RGBColor(180,  0,   0),
    "lt_red":      RGBColor(255,220, 220),
    "orange":      RGBColor(180, 80,   0),
    "lt_orange":   RGBColor(255,235, 200),
    "teal":        RGBColor(0,  110, 110),
    "lt_teal":     RGBColor(200,240, 238),
    "dark_gray":   RGBColor(50,  50,  50),
    "mid_gray":    RGBColor(120,120, 120),
    "lt_gray":     RGBColor(235,235, 240),
    "amber":       RGBColor(200,130,   0),
    "lt_amber":    RGBColor(255,243, 200),
}

def new_slide():
    return prs.slides.add_slide(BLANK)

# ── Low-level helpers ──────────────────────────────────────────────────────
def bg(slide, color_key):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C[color_key]

def rect(slide, l, t, w, h, fill_key, line_key=None, line_w=Pt(0)):
    shp = slide.shapes.add_shape(1,
        I(l), I(t), I(w), I(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = C[fill_key]
    if line_key:
        shp.line.color.rgb = C[line_key]
        shp.line.width = line_w
    else:
        shp.line.fill.background()
    return shp

def roundrect(slide, l, t, w, h, fill_key, line_key=None, line_w=Pt(1.5), radius=0.05):
    shp = slide.shapes.add_shape(5,   # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE = 5
        I(l), I(t), I(w), I(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = C[fill_key]
    if line_key:
        shp.line.color.rgb = C[line_key]
        shp.line.width = line_w
    else:
        shp.line.fill.background()
    # set corner radius
    adj = shp.adjustments
    if adj:
        adj[0] = radius
    return shp

def label(slide, text, l, t, w, h, size=11, bold=False, color_key="white",
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = C[color_key]
    return box

def label2(slide, line1, line2, l, t, w, h, size1=11, size2=9,
           bold1=True, color1="white", color2="sky", align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = align
    r1 = p1.add_run()
    r1.text = line1
    r1.font.size = Pt(size1)
    r1.font.bold = bold1
    r1.font.color.rgb = C[color1]
    p2 = tf.add_paragraph()
    p2.alignment = align
    r2 = p2.add_run()
    r2.text = line2
    r2.font.size = Pt(size2)
    r2.font.italic = True
    r2.font.color.rgb = C[color2]
    return box

def arrow(slide, x1, y1, x2, y2, color_key="mid_gray", w=Pt(1.5)):
    """Draw a simple line arrow from (x1,y1) to (x2,y2) in inches."""
    connector = slide.shapes.add_connector(
        1,  # straight
        I(x1), I(y1), I(x2), I(y2)
    )
    connector.line.color.rgb = C[color_key]
    connector.line.width = w
    # add arrow head
    ln = connector.line._ln
    tail_end = etree.SubElement(ln, qn('a:tailEnd'))
    tail_end.set('type', 'none')
    head_end = etree.SubElement(ln, qn('a:headEnd'))
    head_end.set('type', 'triangle')
    head_end.set('w', 'med')
    head_end.set('len', 'med')
    return connector

def darrow(slide, x1, y1, x2, y2, color_key="blue", w=Pt(2)):
    """Double-headed arrow."""
    connector = slide.shapes.add_connector(1,
        I(x1), I(y1), I(x2), I(y2))
    connector.line.color.rgb = C[color_key]
    connector.line.width = w
    ln = connector.line._ln
    for end_tag in ['a:tailEnd', 'a:headEnd']:
        e = etree.SubElement(ln, qn(end_tag))
        e.set('type', 'triangle')
        e.set('w', 'med')
        e.set('len', 'med')
    return connector

def hdr(slide, title, subtitle=""):
    rect(slide, 0, 0, 13.33, 1.15, "navy")
    rect(slide, 0, 1.15, 13.33, 0.05, "gold")
    label(slide, title, 0.3, 0.08, 12.0, 0.7, size=24, bold=True,
          color_key="white", align=PP_ALIGN.LEFT)
    if subtitle:
        label(slide, subtitle, 0.3, 0.75, 12.0, 0.38, size=11,
              italic=True, color_key="sky", align=PP_ALIGN.LEFT)

def chip(slide, text, l, t, w, h, fill_key, text_color="white", size=9, bold=True):
    r = roundrect(slide, l, t, w, h, fill_key, line_key=None)
    label(slide, text, l+0.04, t+0.03, w-0.08, h-0.06,
          size=size, bold=bold, color_key=text_color, align=PP_ALIGN.CENTER)

def divider(slide, y, color_key="sky"):
    rect(slide, 0.3, y, 12.73, 0.015, color_key)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════════
def slide_cover():
    s = new_slide()
    bg(s, "dark_navy")

    # Background gradient strips
    rect(s, 0, 0, 13.33, 7.5, "dark_navy")
    rect(s, 0, 0, 0.6, 7.5, "navy")
    rect(s, 12.73, 0, 0.6, 7.5, "navy")

    # Gold accent bars
    rect(s, 0, 2.7, 13.33, 0.04, "gold")
    rect(s, 0, 4.85, 13.33, 0.04, "gold")

    # Title block
    label(s, "🏦", 6.1, 0.5, 1.2, 1.0, size=36, align=PP_ALIGN.CENTER, color_key="white")
    label(s, "Internal Bank Employee Assistant",
          0.8, 1.4, 11.7, 1.1, size=32, bold=True, color_key="white",
          align=PP_ALIGN.CENTER)
    label(s, "Architecture Diagram",
          1.5, 2.35, 10.3, 0.7, size=22, bold=True, color_key="gold",
          align=PP_ALIGN.CENTER)

    # Subtitle block
    label(s, "Multi-Agent AI System  |  NexaBank Internal Division  |  Capstone v2.0",
          1.5, 2.88, 10.3, 0.55, size=13, italic=True,
          color_key="sky", align=PP_ALIGN.CENTER)

    # Technology row
    chips_data = [
        ("Python 3.10+", "blue"),
        ("Streamlit UI", "green"),
        ("7 Subagents", "purple"),
        ("3 Skill Modules", "teal"),
        ("6 Lifecycle Hooks", "orange"),
        ("Offline / No API", "red"),
    ]
    x = 0.65
    for ch, col in chips_data:
        chip(s, ch, x, 3.65, 1.85, 0.44, col, "white", 10)
        x += 1.97

    # Slide list
    slides_info = [
        ("Slide 2", "Full System Architecture — 6-layer stack overview"),
        ("Slide 3", "Agent Layer — OrchestratorAgent routing decision tree"),
        ("Slide 4", "End-to-End Data Flow — Single query through all 9 pipeline steps"),
        ("Slide 5", "Skills × Agents Matrix + Hook Lifecycle timeline"),
        ("Slide 6", "Self-Healing Workflow — FallbackAgent auto-chain deep-dive"),
    ]
    label(s, "What's Inside:", 1.5, 4.3, 10.3, 0.38, size=11, bold=True,
          color_key="gold", align=PP_ALIGN.LEFT)
    for i, (slide_num, desc) in enumerate(slides_info):
        label(s, f"  {slide_num}:  {desc}",
              1.5, 4.7 + i*0.47, 10.3, 0.42, size=10,
              color_key="sky", align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — FULL SYSTEM ARCHITECTURE (6-layer stack)
# ════════════════════════════════════════════════════════════════════════════
def slide_full_arch():
    s = new_slide()
    bg(s, "lt_gray")
    hdr(s, "Full System Architecture",
        "6-Layer Stack  |  Streamlit UI → Orchestrator → Agents → Skills → Knowledge Base → Audit Logs")

    # Layer definitions: (label, y, height, bg_key, border_key, text_color)
    layers = [
        ("LAYER 1 — PRESENTATION", 1.28, 0.96, "sky",      "blue",   "navy"),
        ("LAYER 2 — ORCHESTRATION", 2.31, 0.72, "navy",    "dark_navy","white"),
        ("LAYER 3 — AGENTS",       3.10, 1.55, "lt_purple","purple", "dark_gray"),
        ("LAYER 4 — SKILLS",       4.72, 0.82, "lt_green", "green",  "dark_gray"),
        ("LAYER 5 — KNOWLEDGE BASE",5.61, 0.72,"lt_orange","orange", "dark_gray"),
        ("LAYER 6 — AUDIT LOGS",   6.40, 0.88, "lt_teal",  "teal",  "dark_gray"),
    ]
    for (lname, y, h, bg_k, bdr, tc) in layers:
        # Layer band
        rect(s, 0.18, y, 12.97, h, bg_k)
        # Left label strip
        rect(s, 0.18, y, 1.55, h, bdr)
        label(s, lname, 0.21, y + h/2 - 0.15, 1.48, 0.35,
              size=7, bold=True, color_key="white",
              align=PP_ALIGN.CENTER)

    # ── LAYER 1 content — Streamlit app.py ────────────────────────────────
    roundrect(s, 1.9, 1.32, 3.2, 0.82, "blue", "dark_navy", Pt(2))
    label2(s, "app.py", "Streamlit Chat UI  |  Session State  |  st.cache_resource",
           1.95, 1.36, 3.1, 0.74, size1=13, size2=9, color1="white", color2="sky")
    roundrect(s, 5.3, 1.32, 2.5, 0.82, "green", "dark_navy", Pt(2))
    label2(s, "SVG Logos + Badges", "Sidebar  |  Agent Chips  |  Policy Badges",
           5.35, 1.36, 2.4, 0.74, size1=11, size2=9, color1="white", color2="lt_green")
    roundrect(s, 8.0, 1.32, 3.3, 0.82, "teal", "dark_navy", Pt(2))
    label2(s, "Audit Log Viewer", "read_audit_log()  |  Last N lines  |  3 log files",
           8.05, 1.36, 3.2, 0.74, size1=11, size2=9, color1="white", color2="lt_teal")

    # ── LAYER 2 content — Orchestrator ────────────────────────────────────
    roundrect(s, 1.9, 2.35, 10.5, 0.62, "dark_navy", "gold", Pt(2))
    label(s, "OrchestratorAgent  —  Session Coordinator  |  Routing Engine  |  Hook Dispatcher  |  Self-Healing Controller",
          2.0, 2.38, 10.3, 0.55, size=11, bold=True, color_key="gold",
          align=PP_ALIGN.CENTER)

    # ── LAYER 3 content — 7 Agents ────────────────────────────────────────
    agents = [
        ("QueryClassifier\nAgent", "SubAgent 1", "green",   1.85),
        ("PolicySearch\nAgent",    "SubAgent 2", "blue",    3.37),
        ("MultiPolicy\nAgent",     "SubAgent 3", "purple",  4.89),
        ("Summary\nAgent",         "SubAgent 4", "purple",  6.41),
        ("Escalation\nAgent",      "SubAgent 5", "red",     7.93),
        ("Fallback\nAgent",        "SubAgent 6", "orange",  9.45),
        ("AuditLogger\nAgent",     "SubAgent 7", "teal",   10.97),
    ]
    for (name, sub_label, col, lx) in agents:
        roundrect(s, lx, 3.16, 1.4, 1.42, col, "dark_navy", Pt(1))
        label(s, name, lx+0.06, 3.2, 1.28, 0.72,
              size=9, bold=True, color_key="white", align=PP_ALIGN.CENTER)
        label(s, sub_label, lx+0.06, 3.95, 1.28, 0.28,
              size=8, italic=True, color_key="sky", align=PP_ALIGN.CENTER)
        # arrow from orchestrator down to agent
        arrow(s, lx+0.7, 2.97, lx+0.7, 3.16, "mid_gray", Pt(1.2))

    # ── LAYER 4 content — 3 Skills ────────────────────────────────────────
    skills = [
        ("search_skill.py",  "search_policy()  •  extract_keywords()  •  cross_reference()", "green",  1.9),
        ("format_skill.py",  "POLICY_DISPLAY_NAMES  •  format_answer()  •  summarize_policy()  •  generate_escalation_email()", "purple", 5.2),
        ("audit_skill.py",   "log_query()  •  log_unresolved()  •  log_session_event()  •  detect_compliance_risk()  •  read_audit_log()", "teal", 8.5),
    ]
    for (sname, fns, col, lx) in skills:
        roundrect(s, lx, 4.76, 3.05, 0.70, "white", col, Pt(2))
        label(s, sname, lx+0.08, 4.79, 2.9, 0.28, size=10, bold=True,
              color_key="dark_gray", align=PP_ALIGN.LEFT)
        label(s, fns, lx+0.08, 5.06, 2.9, 0.33, size=8, italic=True,
              color_key="mid_gray", align=PP_ALIGN.LEFT)

    # ── LAYER 5 content — Knowledge Base ──────────────────────────────────
    policies = [
        ("loan_policy.txt",              "LP-001", "blue"),
        ("kyc_policy.txt",               "KYC-002","purple"),
        ("customer_complaint_policy.txt","CCP-003","red"),
        ("credit_card_policy.txt",       "CCP-004","green"),
        ("account_opening_policy.txt",   "AOP-005","teal"),
    ]
    xp = 1.9
    for (fname, code, col) in policies:
        roundrect(s, xp, 5.65, 2.18, 0.60, "white", col, Pt(1.5))
        label(s, fname, xp+0.07, 5.67, 2.04, 0.28, size=8, bold=True,
              color_key="dark_gray", align=PP_ALIGN.CENTER)
        label(s, code,  xp+0.07, 5.92, 2.04, 0.22, size=8, italic=True,
              color_key="mid_gray", align=PP_ALIGN.CENTER)
        xp += 2.24
    roundrect(s, 1.9, 5.65, 0.0, 0.0, "lt_orange")  # spacer — already handled

    # utils/ note
    roundrect(s, 11.3, 5.65, 1.7, 0.60, "lt_orange", "orange", Pt(1.5))
    label(s, "utils/\nKEYWORD_MAP", 11.35, 5.67, 1.6, 0.55, size=8, bold=True,
          color_key="orange", align=PP_ALIGN.CENTER)

    # ── LAYER 6 content — Logs ────────────────────────────────────────────
    logs = [
        ("query_audit.log",        "Every query  |  MATCHED / FALLBACK", "teal"),
        ("unresolved_queries.log", "found=False  |  escalation_hook",     "orange"),
        ("session_events.log",     "All hooks  |  Timing  |  Alerts",     "blue"),
    ]
    xl = 1.9
    for (lname, detail, col) in logs:
        roundrect(s, xl, 6.44, 3.55, 0.74, "white", col, Pt(2))
        label(s, lname,  xl+0.1, 6.46, 3.35, 0.3, size=9, bold=True,
              color_key="dark_gray", align=PP_ALIGN.LEFT)
        label(s, detail, xl+0.1, 6.73, 3.35, 0.3, size=8, italic=True,
              color_key="mid_gray", align=PP_ALIGN.LEFT)
        xl += 3.72

    # Right side: hooks.py
    roundrect(s, 11.3, 6.44, 1.7, 0.74, "lt_purple", "purple", Pt(1.5))
    label(s, "hooks.py\n6 Hooks", 11.35, 6.46, 1.6, 0.68, size=9, bold=True,
          color_key="purple", align=PP_ALIGN.CENTER)

    # Down arrows: agents → skills
    arrow(s, 3.2, 4.58, 3.2, 4.76, "green", Pt(1.2))
    arrow(s, 6.7, 4.58, 6.7, 4.76, "purple", Pt(1.2))
    arrow(s, 11.72, 4.58, 11.72, 4.76, "teal", Pt(1.2))

    # Down arrows: skills → knowledge base
    arrow(s, 3.4, 5.46, 3.4, 5.65, "green", Pt(1.2))
    arrow(s, 6.7, 5.46, 6.7, 5.65, "purple", Pt(1.2))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — AGENT LAYER: ROUTING & DECISION TREE
# ════════════════════════════════════════════════════════════════════════════
def slide_routing():
    s = new_slide()
    bg(s, "white")
    hdr(s, "Agent Layer — OrchestratorAgent Routing & Decision Tree",
        "QueryClassifierAgent classifies intent → OrchestratorAgent applies routing rules → correct subagent invoked")

    # Central orchestrator box
    roundrect(s, 4.2, 1.35, 5.0, 0.82, "dark_navy", "gold", Pt(3))
    label(s, "OrchestratorAgent", 4.3, 1.38, 4.8, 0.42,
          size=15, bold=True, color_key="gold", align=PP_ALIGN.CENTER)
    label(s, "Session Coordinator  |  Routing Engine  |  Self-Healing Controller",
          4.3, 1.76, 4.8, 0.32, size=9, italic=True,
          color_key="sky", align=PP_ALIGN.CENTER)

    # pre_query_hook input arrow
    roundrect(s, 0.25, 1.52, 2.0, 0.52, "blue", "navy", Pt(1.5))
    label(s, "pre_query_hook\nValidates input first", 0.3, 1.53, 1.9, 0.48,
          size=9, bold=True, color_key="white", align=PP_ALIGN.CENTER)
    arrow(s, 2.25, 1.78, 4.2, 1.78, "blue", Pt(2))
    label(s, "valid=True\nquestion", 2.35, 1.58, 1.7, 0.38, size=8,
          color_key="blue", align=PP_ALIGN.CENTER)

    # QueryClassifier above
    roundrect(s, 4.2, 2.32, 5.0, 0.6, "green", "dark_navy", Pt(1.5))
    label(s, "QueryClassifierAgent  (SubAgent 1)  —  Detects 4 Intents",
          4.3, 2.35, 4.8, 0.52, size=10, bold=True, color_key="white",
          align=PP_ALIGN.CENTER)
    arrow(s, 6.7, 2.17, 6.7, 2.32, "green", Pt(2))

    # Four intent branches coming down from orchestrator
    # Layout: summary | escalation | policy_query (single) | policy_query (multi) | out_of_scope
    branches = [
        # (label_top, x_center, agent_name, sub, fill, detail)
        ("summary\nintent",        1.55,  "SummaryAgent",       "SubAgent 4", "purple",
         "POLICY_NAME_HINTS\n→ summarize_policy()"),
        ("escalation\nintent",     4.05,  "EscalationAgent",    "SubAgent 5", "red",
         "ESCALATION_CONTACTS\n→ email draft"),
        ("policy_query\nsingle",   6.55,  "PolicySearchAgent",  "SubAgent 2", "blue",
         "KEYWORD_MAP scoring\n→ best match"),
        ("policy_query\nmulti-domain", 9.1, "MultiPolicyAgent","SubAgent 3", "purple",
         "cross_reference()\ntop≥2 & second≥1"),
        ("out_of_\nscope",        11.7,  "Direct Fallback",    "no subagent","dark_gray",
         "Empty / <3 chars\n→ polite message"),
    ]

    for (intent_lbl, xc, agent_nm, sub_lbl, col, detail) in branches:
        # Intent chip at top
        roundrect(s, xc-0.78, 3.08, 1.56, 0.58, "lt_gray", col, Pt(1.5))
        label(s, intent_lbl, xc-0.75, 3.1, 1.5, 0.52,
              size=9, bold=True, color_key=col, align=PP_ALIGN.CENTER)
        # Down arrow
        arrow(s, xc, 3.66, xc, 3.85, col, Pt(1.8))
        # Agent box
        roundrect(s, xc-0.88, 3.85, 1.76, 0.9, col, "dark_navy", Pt(1.5))
        label(s, agent_nm,  xc-0.84, 3.88, 1.68, 0.44,
              size=9, bold=True, color_key="white", align=PP_ALIGN.CENTER)
        label(s, sub_lbl, xc-0.84, 4.30, 1.68, 0.28,
              size=8, italic=True, color_key="sky", align=PP_ALIGN.CENTER)
        # Detail box
        roundrect(s, xc-0.88, 4.82, 1.76, 0.62, "lt_gray", col, Pt(1))
        label(s, detail, xc-0.84, 4.84, 1.68, 0.56,
              size=8, color_key="dark_gray", align=PP_ALIGN.CENTER)

    # Arrows from orchestrator to intent chips
    orch_cx = 6.7  # center of orchestrator
    for (_, xc, *_) in branches:
        arrow(s, orch_cx, 2.92, xc, 3.08, "mid_gray", Pt(1.2))

    # Self-healing arrow: PolicySearchAgent → FallbackAgent
    roundrect(s, 6.95, 5.6, 1.76, 0.7, "orange", "dark_navy", Pt(2))
    label(s, "FallbackAgent\n(SubAgent 6)", 7.0, 5.63, 1.66, 0.62,
          size=9, bold=True, color_key="white", align=PP_ALIGN.CENTER)
    arrow(s, 7.43, 4.75, 7.43, 5.6, "orange", Pt(2))
    label(s, "found=False\n(self-heal)", 7.5, 5.05, 1.4, 0.42,
          size=8, italic=True, color_key="orange", align=PP_ALIGN.LEFT)

    # AuditLogger always runs
    roundrect(s, 10.95, 5.6, 1.76, 0.7, "teal", "dark_navy", Pt(2))
    label(s, "AuditLogger\nAgent (SubAgent 7)", 11.0, 5.63, 1.66, 0.62,
          size=9, bold=True, color_key="white", align=PP_ALIGN.CENTER)
    label(s, "← runs after EVERY query", 9.1, 5.78, 1.8, 0.32,
          size=8, italic=True, color_key="teal", align=PP_ALIGN.RIGHT)

    # Hook bar at bottom
    rect(s, 0.18, 6.5, 12.97, 0.78, "lt_teal")
    label(s, "🪝  Hooks (auto-fired — not routed by Orchestrator):",
          0.3, 6.52, 3.5, 0.3, size=9, bold=True, color_key="teal")
    hook_chips = [
        ("pre_query_hook",  "blue"),    ("post_answer_hook", "blue"),
        ("session_start",   "green"),   ("session_end",      "green"),
        ("escalation_hook", "orange"),  ("compliance_alert", "red"),
    ]
    xh = 3.9
    for (hname, col) in hook_chips:
        chip(s, hname, xh, 6.6, 1.52, 0.33, col, "white", 8)
        xh += 1.58


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — END-TO-END DATA FLOW
# ════════════════════════════════════════════════════════════════════════════
def slide_data_flow():
    s = new_slide()
    bg(s, "white")
    hdr(s, "End-to-End Data Flow — Single Policy Query",
        "Question: 'What documents are needed for a loan?'  |  9 steps from input to screen")

    steps = [
        # (num, name, role, x, y, fill, arrow_to_x, arrow_to_y)
        ("1", "Employee",            "Types question",                    0.22, 1.42, "blue"),
        ("2", "app.py",              "Calls orchestrator.answer()",       2.28, 1.42, "navy"),
        ("3", "pre_query_hook",      "Validates & logs",                  4.35, 1.42, "green"),
        ("4", "QueryClassifier",     "intent = policy_query",             6.42, 1.42, "green"),
        ("5", "Orchestrator",        "Routes: single domain",             8.48, 1.42, "dark_navy"),
        ("6", "search_skill\n.search_policy()", "Scores KEYWORD_MAP",    10.55, 1.42, "purple"),
        ("7", "PolicySearch\nAgent", "format_answer()",                   0.22, 3.52, "blue"),
        ("8", "AuditLogger\nAgent",  "log_query()",                       2.28, 3.52, "teal"),
        ("9", "app.py UI",           "Render result",                     4.35, 3.52, "navy"),
    ]

    bw, bh = 2.0, 1.2
    for i, (num, name, role, lx, ly, col) in enumerate(steps):
        roundrect(s, lx, ly, bw, bh, col, "dark_navy", Pt(2))
        # Step number circle
        roundrect(s, lx+0.05, ly+0.05, 0.36, 0.36, "gold", "dark_navy", Pt(0))
        label(s, num, lx+0.07, ly+0.07, 0.32, 0.3, size=11, bold=True,
              color_key="dark_navy", align=PP_ALIGN.CENTER)
        label(s, name, lx+0.45, ly+0.08, 1.5, 0.52,
              size=10, bold=True, color_key="white", align=PP_ALIGN.LEFT)
        label(s, role, lx+0.08, ly+0.65, 1.85, 0.48,
              size=9, italic=True, color_key="sky", align=PP_ALIGN.LEFT)

    # Arrows row 1 (steps 1-6, horizontal)
    for i in range(5):
        lx1 = steps[i][3] + bw
        ly1 = steps[i][4] + bh/2
        lx2 = steps[i+1][3]
        arrow(s, lx1, ly1, lx2, ly1, "mid_gray", Pt(1.8))

    # Turn arrow from step 6 down to row 2 step 7
    arrow(s, 11.55, 2.02, 11.55, 3.28, "mid_gray", Pt(1.8))
    arrow(s, 11.55, 3.28, 2.22, 3.28, "mid_gray", Pt(1.8))
    arrow(s, 2.22, 3.28, 2.22 + 0.0, 3.28 + 0.24, "mid_gray", Pt(1.8))
    # row 2 arrows (steps 7-9)
    for i in [6, 7]:
        lx1 = steps[i][3] + bw
        ly1 = steps[i][4] + bh/2
        lx2 = steps[i+1][3]
        arrow(s, lx1, ly1, lx2, ly1, "mid_gray", Pt(1.8))

    # Data flowing between steps
    data_labels = [
        (1.32, 1.78, "question str"),
        (3.38, 1.78, "question"),
        (5.45, 1.78, "{valid, q}"),
        (7.51, 1.78, "{intent}"),
        (9.57, 1.78, "question + policies"),
        (10.2, 3.14, "{matched_key, snippet, score}"),
        (3.38, 3.9,  "{answer, found}"),
        (5.45, 3.9,  "{logged}"),
    ]
    for (lx, ly, txt) in data_labels:
        label(s, txt, lx, ly, 1.85, 0.28, size=7, italic=True,
              color_key="mid_gray", align=PP_ALIGN.CENTER)

    # Side hooks column
    rect(s, 10.7, 3.36, 2.5, 3.9, "lt_teal")
    label(s, "🪝  Hooks Also Fire:", 10.78, 3.42, 2.3, 0.33,
          size=10, bold=True, color_key="teal")
    hooks_side = [
        ("compliance_alert_hook",  "AuditLogger risk=True → alert str",       "red"),
        ("escalation_hook",        "found=False AND NOT escalated → log",     "orange"),
        ("post_answer_hook",       "Response time + outcome → session log",   "blue"),
    ]
    yh = 3.85
    for (hn, hd, hc) in hooks_side:
        roundrect(s, 10.78, yh, 2.3, 0.78, "white", hc, Pt(1.5))
        label(s, hn, 10.83, yh+0.04, 2.2, 0.3, size=9, bold=True,
              color_key=hc, align=PP_ALIGN.LEFT)
        label(s, hd, 10.83, yh+0.36, 2.2, 0.36, size=8, italic=True,
              color_key="dark_gray", align=PP_ALIGN.LEFT)
        yh += 0.89

    # Result display box
    roundrect(s, 0.22, 5.0, 10.2, 1.2, "sky", "navy", Pt(2))
    label(s, "✅  Result delivered to app.py → UI renders:", 0.35, 5.04, 9.8, 0.32,
          size=10, bold=True, color_key="navy")
    result_items = [
        ("📚 Loan Policy (LP-001)", "policy badge"),
        ("🤖 PolicySearchAgent | policy_query", "agent chip"),
        ("📝 Verbatim answer text", "green bot bubble"),
        ("🔒 No compliance alert", "risk=False"),
        ("📋 Audit logged", "query_audit.log"),
    ]
    xi = 0.35
    for (ri, rd) in result_items:
        roundrect(s, xi, 5.44, 1.96, 0.65, "white", "blue", Pt(1))
        label(s, ri, xi+0.06, 5.47, 1.84, 0.28, size=8, bold=True,
              color_key="navy", align=PP_ALIGN.CENTER)
        label(s, rd, xi+0.06, 5.73, 1.84, 0.28, size=7, italic=True,
              color_key="mid_gray", align=PP_ALIGN.CENTER)
        xi += 2.04


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SKILLS × AGENTS MATRIX + HOOK LIFECYCLE
# ════════════════════════════════════════════════════════════════════════════
def slide_skills_hooks():
    s = new_slide()
    bg(s, "white")
    hdr(s, "Skills × Agents Usage Matrix  +  Hook Lifecycle Timeline",
        "Which agent calls which skill  |  When each hook fires in the pipeline")

    # ── LEFT: Skills × Agents matrix ──────────────────────────────────────
    label(s, "Skills × Agents Matrix", 0.22, 1.28, 6.5, 0.35,
          size=12, bold=True, color_key="navy")

    # Agent column labels (top row)
    agents_m = ["QCAgent", "PSAgent", "MPAgent", "SummAgent", "EscAgent", "FbackAgent", "ALogAgent"]
    agent_colors = ["green","blue","purple","purple","red","orange","teal"]
    col_w = 0.78
    x0 = 1.75
    for i, (ag, col) in enumerate(zip(agents_m, agent_colors)):
        roundrect(s, x0 + i*col_w, 1.66, col_w-0.05, 0.45, col, "dark_navy", Pt(1))
        label(s, ag, x0 + i*col_w + 0.02, 1.68, col_w-0.09, 0.38,
              size=7, bold=True, color_key="white", align=PP_ALIGN.CENTER)

    # Skill row labels + matrix cells
    skills_m = [
        ("search_skill\n.search_policy()",   ["", "✓", "✓", "✓*", "✓", "", ""],      "green"),
        ("search_skill\n.extract_keywords()", ["✓", "", "", "", "", "", ""],           "green"),
        ("search_skill\n.cross_reference()", ["", "", "✓", "", "", "✓", ""],          "green"),
        ("format_skill\n.format_answer()",   ["", "✓", "✓", "", "", "", ""],          "purple"),
        ("format_skill\n.summarize_policy()", ["", "", "", "✓", "", "", ""],          "purple"),
        ("format_skill\ngenerate_email()",   ["", "", "", "", "✓", "", ""],            "purple"),
        ("format_skill\nPOLICY_DISPLAY_NAMES",["", "✓", "✓", "✓", "✓", "✓", ""],    "purple"),
        ("audit_skill\n.log_query()",        ["", "", "", "", "", "", "✓"],            "teal"),
        ("audit_skill\n.log_unresolved()",   ["", "", "", "", "", "", "✓"],            "teal"),
        ("audit_skill\n.log_session_event()",["", "", "", "", "", "", "✓"],            "teal"),
        ("audit_skill\n.detect_risk()",      ["", "", "", "", "", "", "✓"],            "teal"),
        ("audit_skill\n.read_audit_log()",   ["", "", "", "", "", "", ""],             "teal"),
    ]
    y0 = 2.16
    row_h = 0.38
    for ri, (sk_name, cells, sk_col) in enumerate(skills_m):
        ry = y0 + ri*row_h
        # skill label
        roundrect(s, 0.22, ry, 1.48, row_h-0.04, "lt_gray", sk_col, Pt(1))
        label(s, sk_name, 0.24, ry+0.02, 1.44, row_h-0.08,
              size=7, color_key="dark_gray", align=PP_ALIGN.LEFT)
        # cells
        for ci, cell in enumerate(cells):
            cx = x0 + ci*col_w
            fill_c = agent_colors[ci] if cell == "✓" else "white"
            roundrect(s, cx, ry, col_w-0.05, row_h-0.04,
                      fill_c if cell == "✓" else "white",
                      agent_colors[ci], Pt(0.5))
            if cell:
                label(s, cell, cx, ry+0.04, col_w-0.05, row_h-0.12,
                      size=10, bold=True,
                      color_key="white" if cell == "✓" else agent_colors[ci],
                      align=PP_ALIGN.CENTER)

    label(s, "* SummaryAgent uses search_policy() only as fallback (when no POLICY_NAME_HINT matches)",
          0.22, y0 + len(skills_m)*row_h + 0.05, 6.5, 0.28, size=7, italic=True,
          color_key="mid_gray")

    # ── RIGHT: Hook lifecycle timeline ────────────────────────────────────
    divider(s, 1.28)
    rect(s, 7.0, 1.28, 0.03, 6.1, "sky")
    label(s, "Hook Lifecycle — When Each Hook Fires",
          7.1, 1.28, 5.9, 0.35, size=12, bold=True, color_key="navy")

    pipeline_steps = [
        "① Employee types question",
        "② app.py calls orchestrator",
        "③ Pipeline step: pre validation",
        "④ QueryClassifierAgent",
        "⑤ SubAgent executes",
        "⑥ AuditLoggerAgent runs",
        "⑦ Compliance / Escalation check",
        "⑧ post answer logging",
        "⑨ Result → UI",
    ]
    hooks_hl = [
        # (hook_name, fires_at_step_index, color, note)
        ("pre_query_hook",        2, "blue",   "valid=True/False → { valid, question }"),
        ("session_start_hook",    1, "green",  "Once: OrchestratorAgent.__init__()"),
        ("compliance_alert_hook", 5, "red",    "IF detect_compliance_risk() = True"),
        ("escalation_hook",       6, "orange", "IF found=False AND escalated=False"),
        ("post_answer_hook",      7, "blue",   "Every query — logs timing + outcome"),
        ("session_end_hook",      8, "green",  "Employee clicks Clear Chat"),
    ]

    tl_x = 7.18
    tl_y0 = 1.72
    step_h = 0.57

    for i, step in enumerate(pipeline_steps):
        sy = tl_y0 + i*step_h
        # step bubble
        roundrect(s, tl_x, sy, 2.2, 0.45, "lt_gray", "navy", Pt(1))
        label(s, step, tl_x+0.06, sy+0.05, 2.08, 0.34, size=9,
              color_key="dark_gray", align=PP_ALIGN.LEFT)
        # connector line down
        if i < len(pipeline_steps)-1:
            rect(s, tl_x+1.08, sy+0.45, 0.04, step_h-0.45, "mid_gray")

    # Hook arrows pointing in from right
    for (hname, step_idx, col, note) in hooks_hl:
        hy = tl_y0 + step_idx*step_h + 0.22
        # arrow from right
        roundrect(s, 10.0, hy-0.22, 3.15, 0.44, "white", col, Pt(1.5))
        label(s, hname, 10.06, hy-0.2, 2.0, 0.24, size=9, bold=True,
              color_key=col, align=PP_ALIGN.LEFT)
        label(s, note, 10.06, hy+0.04, 3.03, 0.2, size=7, italic=True,
              color_key="mid_gray", align=PP_ALIGN.LEFT)
        arrow(s, 10.0, hy, tl_x+2.2, hy, col, Pt(1.8))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SELF-HEALING WORKFLOW
# ════════════════════════════════════════════════════════════════════════════
def slide_selfheal():
    s = new_slide()
    bg(s, "white")
    hdr(s, "Self-Healing Workflow — FallbackAgent Auto-Chain",
        "Question: 'What is the HR leave policy?'  |  PolicySearchAgent fails → Orchestrator auto-invokes FallbackAgent")

    # Heading note
    roundrect(s, 0.22, 1.28, 12.89, 0.5, "lt_amber", "amber", Pt(2))
    label(s, "⚠️  Key Concept: Self-healing is invisible to the employee. "
           "PolicySearchAgent returning found=False is the ONLY trigger. "
           "EscalationAgent is NOT part of this chain — it is a separate agent for explicit 'escalate' intent.",
          0.35, 1.3, 12.65, 0.44, size=9, italic=True, color_key="orange")

    # Step boxes — two rows
    steps_sh = [
        # (num, title, detail, x, y, fill)
        ("1", "Employee\nTypes Question",
         "HR leave policy for\nannual leave",
         0.22, 2.0, "blue"),
        ("2", "pre_query_hook",
         "Length > 3 ✓\nLogs INCOMING",
         2.85, 2.0, "blue"),
        ("3", "QueryClassifier\nAgent",
         "No summary/escalation\ntriggers → policy_query",
         5.48, 2.0, "green"),
        ("4", "Orchestrator\nRouting",
         "cross_reference all 5\npolicies → all score=0",
         8.12, 2.0, "dark_navy"),
        ("5", "PolicySearch\nAgent",
         "KEYWORD_MAP search\n→ score=0 all 5 → found=False",
         10.75, 2.0, "blue"),
        ("6", "Orchestrator\nSelf-Heal Trigger",
         "Detects found=False\n→ auto-invokes FallbackAgent",
         0.22, 4.1, "dark_navy"),
        ("7", "FallbackAgent\n(SubAgent 6)",
         "cross_reference at\nscore>0 → also 0 → no partial",
         2.85, 4.1, "orange"),
        ("8", "escalation_hook",
         "found=False AND\nescalated=False → FIRES",
         5.48, 4.1, "red"),
        ("9", "AuditLogger\nAgent",
         "query_audit.log:\nSTATUS=FALLBACK\nunresolved_queries.log: logged",
         8.12, 4.1, "teal"),
        ("10","app.py UI\nRenders",
         "Amber bubble\nFallbackAgent chip\nContact details",
         10.75, 4.1, "navy"),
    ]

    bw2, bh2 = 2.4, 1.88

    for (num, title, detail, lx, ly, col) in steps_sh:
        roundrect(s, lx, ly, bw2, bh2, col, "dark_navy", Pt(2))
        # number badge
        roundrect(s, lx+0.06, ly+0.06, 0.4, 0.4, "gold", "dark_navy", Pt(0))
        label(s, num, lx+0.08, ly+0.08, 0.36, 0.32,
              size=11, bold=True, color_key="dark_navy", align=PP_ALIGN.CENTER)
        label(s, title, lx+0.5, ly+0.08, 1.84, 0.52,
              size=10, bold=True, color_key="white", align=PP_ALIGN.LEFT)
        label(s, detail, lx+0.1, ly+0.68, 2.2, 1.1,
              size=9, italic=True, color_key="sky", align=PP_ALIGN.LEFT)

    # Row 1 horizontal arrows (steps 1-5)
    for i in range(4):
        lx1 = steps_sh[i][3] + bw2
        ly1 = steps_sh[i][4] + bh2/2
        lx2 = steps_sh[i+1][3]
        arrow(s, lx1, ly1, lx2, ly1, "mid_gray", Pt(2))

    # Bend from step 5 → step 6 (down then left)
    arrow(s, 12.0, 2.94, 12.0, 3.82, "orange", Pt(2.5))
    arrow(s, 12.0, 3.82, 2.62, 3.82, "orange", Pt(2.5))
    arrow(s, 2.62, 3.82, 2.62, 4.1, "orange", Pt(2.5))
    label(s, "found=False → SELF-HEAL", 5.5, 3.58, 3.8, 0.32,
          size=9, bold=True, color_key="orange", align=PP_ALIGN.CENTER)

    # Row 2 horizontal arrows (steps 6-10)
    for i in range(5, 9):
        lx1 = steps_sh[i][3] + bw2
        ly1 = steps_sh[i][4] + bh2/2
        lx2 = steps_sh[i+1][3]
        arrow(s, lx1, ly1, lx2, ly1, "mid_gray", Pt(2))

    # What employee sees
    roundrect(s, 0.22, 6.2, 12.89, 1.1, "lt_amber", "orange", Pt(2))
    label(s, "👁  What the Employee Sees:",
          0.38, 6.22, 3.0, 0.33, size=10, bold=True, color_key="orange")
    ee_items = [
        "Amber warning bubble (not green)",
        "FallbackAgent chip shown",
        "Contact details: compliance@bank.internal",
        "No 'PolicySearchAgent failed' message — invisible",
    ]
    xi = 0.38
    for item in ee_items:
        label(s, f"• {item}", xi, 6.56, 3.1, 0.55, size=9, color_key="dark_gray")
        xi += 3.16


# ════════════════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════════════════
print("Building Architecture Diagram PPT...")
slide_cover()
print("  ✓ Slide 1 — Cover")
slide_full_arch()
print("  ✓ Slide 2 — Full System Architecture")
slide_routing()
print("  ✓ Slide 3 — Agent Routing & Decision Tree")
slide_data_flow()
print("  ✓ Slide 4 — End-to-End Data Flow")
slide_skills_hooks()
print("  ✓ Slide 5 — Skills × Agents Matrix + Hook Lifecycle")
slide_selfheal()
print("  ✓ Slide 6 — Self-Healing Workflow")

prs.save(OUT)
print(f"\n✅  Saved: {OUT}")
