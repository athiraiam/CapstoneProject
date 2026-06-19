"""
generate_pptx.py
Generates the Presentation Deck (.pptx) for the Bank Employee Assistant capstone.
Run: python generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

# ── Colors ────────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0,  51,  102)
MID_BLUE   = RGBColor(0,  102, 204)
LIGHT_BLUE = RGBColor(204, 229, 255)
WHITE      = RGBColor(255, 255, 255)
GOLD       = RGBColor(255, 193,   7)
DARK_GRAY  = RGBColor(50,  50,  50)
GREEN      = RGBColor(0,  153,  76)

TODAY = datetime.date.today().strftime("%B %d, %Y")


# ── Helpers ───────────────────────────────────────────────────────────────────
def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def bullet_box(slide, title, bullets, left, top, width, height,
               title_color=DARK_BLUE, bullet_color=DARK_GRAY):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    # Title line
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = title_color

    for bullet in bullets:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"  •  {bullet}"
        run.font.size = Pt(13)
        run.font.color.rgb = bullet_color
    return txBox


prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]   # completely blank


# ── SLIDE 1 — Title Slide ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BLUE)

add_rect(s, 0, 2.8, 13.33, 0.06, GOLD)   # gold divider line

add_textbox(s, "🏦  Internal Bank Employee Assistant",
            0.5, 0.8, 12, 1.5, font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s, "AI-Powered Policy Chatbot — Capstone Project",
            0.5, 2.2, 12, 0.8, font_size=20, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_textbox(s, f"Presented by: Bank Technology Division     |     {TODAY}",
            0.5, 3.2, 12, 0.6, font_size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_textbox(s, "Built with  Python  •  Streamlit  •  Local Knowledge Base",
            0.5, 4.0, 12, 0.6, font_size=13, color=GOLD, align=PP_ALIGN.CENTER)


# ── SLIDE 2 — Business Problem ───────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 1.2, DARK_BLUE)
add_textbox(s, "Business Problem", 0.4, 0.2, 12, 0.8, font_size=28, bold=True, color=WHITE)

problems = [
    "Bank employees must reference 5+ internal policies daily",
    "Manual searches are slow and inconsistent",
    "Policy documents are long — finding the right section takes minutes",
    "Inconsistent answers lead to compliance risks",
    "New employees struggle to navigate complex policy documents",
]
for i, prob in enumerate(problems):
    top = 1.4 + i * 0.9
    add_rect(s, 0.4, top, 12.3, 0.75, LIGHT_BLUE)
    add_textbox(s, f"⚠  {prob}", 0.5, top + 0.12, 12, 0.6, font_size=14, color=DARK_BLUE)


# ── SLIDE 3 — Solution Overview ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 1.2, MID_BLUE)
add_textbox(s, "Solution Overview", 0.4, 0.2, 12, 0.8, font_size=28, bold=True, color=WHITE)

add_textbox(s, "An offline AI chatbot that answers policy questions instantly from an internal knowledge base.",
            0.4, 1.3, 12.5, 0.6, font_size=15, color=DARK_BLUE, italic=True)

features = [
    ("📚 Local Knowledge Base", "5 policy documents stored as plain text — no cloud needed"),
    ("🔍 Smart Keyword Search", "Scores all policies simultaneously, picks best match"),
    ("✂️  Snippet Extraction", "Returns only the relevant paragraph, not the full document"),
    ("⚠️  Fallback Safety", "Directs to compliance team when no answer found"),
    ("💬 Chat Interface", "Streamlit chatbot UI — familiar and easy to use"),
]
for i, (title, desc) in enumerate(features):
    left = 0.4 + (i % 3) * 4.3
    top  = 2.1 + (i // 3) * 1.8
    add_rect(s, left, top, 4.0, 1.5, DARK_BLUE)
    add_textbox(s, title, left + 0.1, top + 0.1, 3.8, 0.5, font_size=13, bold=True, color=GOLD)
    add_textbox(s, desc,  left + 0.1, top + 0.6, 3.8, 0.8, font_size=11, color=WHITE)


# ── SLIDE 4 — Agent Architecture ─────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 1.2, DARK_BLUE)
add_textbox(s, "Agent Architecture", 0.4, 0.2, 12, 0.8, font_size=28, bold=True, color=WHITE)

# Draw flow boxes
boxes = [
    (0.3,  2.5, "👤 Employee\nQuestion", MID_BLUE),
    (2.7,  2.5, "🖥️ Streamlit\nChat UI\n(app.py)", MID_BLUE),
    (5.1,  2.5, "🤖 Agent Core\n(agent.py)", DARK_BLUE),
    (7.5,  2.5, "🔍 Keyword\nScorer", MID_BLUE),
    (9.9,  2.5, "📚 Policy\nKnowledge Base", MID_BLUE),
]
for left, top, label, color in boxes:
    add_rect(s, left, top, 2.2, 1.6, color)
    add_textbox(s, label, left + 0.05, top + 0.2, 2.1, 1.2, font_size=12, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

# Arrows (simple text)
for x in [2.55, 4.95, 7.35, 9.75]:
    add_textbox(s, "→", x, 3.0, 0.3, 0.6, font_size=22, bold=True, color=DARK_BLUE)

# Labels below boxes
labels = ["Input", "Display", "Orchestrate", "Score + Extract", "LP/KYC/CC/CCP/AO"]
for i, (left, _, _, _) in enumerate(boxes):
    add_textbox(s, labels[i], left, 4.2, 2.2, 0.5,
                font_size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER, italic=True)

add_textbox(s, "Answer flows back through: Knowledge Base → Agent → UI → Employee",
            0.4, 5.0, 12.5, 0.5, font_size=13, italic=True, color=MID_BLUE, align=PP_ALIGN.CENTER)


# ── SLIDE 5 — Knowledge Base ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 1.2, MID_BLUE)
add_textbox(s, "Knowledge Base — 5 Policy Documents", 0.4, 0.2, 12, 0.8, font_size=26, bold=True, color=WHITE)

policies = [
    ("🏠", "LP-001", "Loan Policy", "Eligibility, types, documents,\ninterest rates, approval process"),
    ("🪪", "KYC-002", "KYC Policy", "Identity verification,\ndocument requirements, renewal"),
    ("📣", "CCP-003", "Complaint Policy", "Channels, timelines,\nescalation levels, compensation"),
    ("💳", "CCP-004", "Credit Card Policy", "Card types, billing, APR,\nlost card, disputes"),
    ("🏦", "AOP-005", "Account Opening", "Account types, documents,\nprocess, dormancy, closure"),
]
for i, (icon, code, name, desc) in enumerate(policies):
    left = 0.4 + i * 2.55
    add_rect(s, left, 1.4, 2.3, 3.5, DARK_BLUE)
    add_textbox(s, icon, left + 0.8, 1.5, 1.0, 0.7, font_size=24, align=PP_ALIGN.CENTER, color=WHITE)
    add_textbox(s, code, left + 0.1, 2.2, 2.1, 0.5, font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(s, name, left + 0.1, 2.7, 2.1, 0.5, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, desc, left + 0.1, 3.2, 2.1, 1.3, font_size=10, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_textbox(s, "All policies stored as plain .txt files — updated by policy owners, no code changes needed.",
            0.4, 5.1, 12.5, 0.5, font_size=12, italic=True, color=MID_BLUE, align=PP_ALIGN.CENTER)


# ── SLIDE 6 — Search & Answer Logic ──────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 1.2, DARK_BLUE)
add_textbox(s, "Skills & Search Logic", 0.4, 0.2, 12, 0.8, font_size=28, bold=True, color=WHITE)

steps = [
    ("1", "Employee asks a question", "e.g. 'What credit score do I need for a credit card?'"),
    ("2", "Keyword Scoring", "All 5 policies scored by keyword overlap — best match selected"),
    ("3", "Snippet Extraction", "Most relevant paragraph extracted (not full document)"),
    ("4", "Answer Returned", "Answer shown in chat with source policy badge"),
    ("5", "No Match → Fallback", "Directs to compliance@bank.internal if no policy matches"),
]
for i, (num, title, detail) in enumerate(steps):
    top = 1.4 + i * 1.0
    add_rect(s, 0.3, top, 0.6, 0.7, MID_BLUE)
    add_textbox(s, num, 0.3, top + 0.1, 0.6, 0.5, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, title,  1.1, top,        5.0, 0.4, font_size=14, bold=True, color=DARK_BLUE)
    add_textbox(s, detail, 1.1, top + 0.38, 11.6, 0.4, font_size=12, color=DARK_GRAY, italic=True)


# ── SLIDE 7 — Governance Framework ───────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 1.2, DARK_BLUE)
add_textbox(s, "Governance Framework", 0.4, 0.2, 12, 0.8, font_size=28, bold=True, color=WHITE)

pillars = [
    ("🔒 Data Privacy", ["No PII collected or stored", "Queries cleared on browser close", "No external API calls"]),
    ("📋 Source Control", ["Every answer cites source policy", "Answers extracted from approved text only", "No AI hallucination"]),
    ("👥 Access Control", ["Internal network only", "Role-based update rights", "Policy owners manage content"]),
    ("🔄 Review Cycle", ["Annual policy review per owner", "Compliance team approves updates", "Version tracked in filenames"]),
]
for i, (title, bullets) in enumerate(pillars):
    left = 0.4 + (i % 2) * 6.4
    top  = 1.5 + (i // 2) * 2.5
    add_rect(s, left, top, 6.0, 2.2, LIGHT_BLUE)
    add_textbox(s, title, left + 0.15, top + 0.1, 5.7, 0.5, font_size=15, bold=True, color=DARK_BLUE)
    for j, b in enumerate(bullets):
        add_textbox(s, f"• {b}", left + 0.2, top + 0.6 + j * 0.45, 5.6, 0.4, font_size=12, color=DARK_GRAY)


# ── SLIDE 8 — Observability & Traceability ───────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 1.2, MID_BLUE)
add_textbox(s, "Observability & Traceability", 0.4, 0.2, 12, 0.8, font_size=28, bold=True, color=WHITE)

items = [
    ("📋 Source Attribution",    "Every answer displays the exact policy name and code (e.g. Loan Policy LP-001)"),
    ("⚠️  Fallback Logging",     "All unmatched queries return an explicit fallback — no silent failures"),
    ("🔍 Answer Traceability",   "Answers are direct text extracts — can be traced back to exact policy line"),
    ("💬 Chat History",          "Full conversation history visible in the session for employee review"),
    ("🚫 No Black Box",          "No AI model — every answer is deterministic and auditable"),
]
for i, (icon_title, desc) in enumerate(items):
    top = 1.5 + i * 1.0
    add_rect(s, 0.3, top, 12.7, 0.8, LIGHT_BLUE if i % 2 == 0 else RGBColor(220, 235, 255))
    add_textbox(s, icon_title, 0.5,  top + 0.15, 3.5, 0.5, font_size=13, bold=True, color=DARK_BLUE)
    add_textbox(s, desc,       4.1,  top + 0.15, 8.7, 0.5, font_size=13, color=DARK_GRAY)


# ── SLIDE 9 — Test Results ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 1.2, DARK_BLUE)
add_textbox(s, "Testing & Evaluation Results", 0.4, 0.2, 12, 0.8, font_size=28, bold=True, color=WHITE)

metrics = [
    ("23", "Total Test Cases"),
    ("23", "Passed"),
    ("0",  "Failed"),
    ("100%", "Pass Rate"),
]
for i, (val, label) in enumerate(metrics):
    left = 0.5 + i * 3.1
    add_rect(s, left, 1.4, 2.8, 1.8, MID_BLUE if i < 3 else GREEN)
    add_textbox(s, val,   left + 0.1, 1.5,  2.6, 1.0, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, label, left + 0.1, 2.5,  2.6, 0.5, font_size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

domain_results = [
    ("Loan Policy",      "4/4",  "100%"),
    ("KYC Policy",       "4/4",  "100%"),
    ("Complaint Policy", "3/3",  "100%"),
    ("Credit Card",      "4/4",  "100%"),
    ("Account Opening",  "3/3",  "100%"),
    ("Fallback Tests",   "5/5",  "100%"),
]
add_textbox(s, "Results by Domain:", 0.4, 3.4, 5, 0.5, font_size=14, bold=True, color=DARK_BLUE)
for i, (domain, score, pct) in enumerate(domain_results):
    left = 0.4 + (i % 3) * 4.2
    top  = 3.9 + (i // 3) * 0.7
    add_textbox(s, f"✅  {domain}: {score} ({pct})", left, top, 4.0, 0.6, font_size=12, color=DARK_GRAY)


# ── SLIDE 10 — Deployment Architecture ───────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 1.2, MID_BLUE)
add_textbox(s, "Deployment Architecture", 0.4, 0.2, 12, 0.8, font_size=28, bold=True, color=WHITE)

options = [
    ("Option A\nLocal Deploy", "• pip install -r requirements.txt\n• streamlit run app.py\n• http://localhost:8501\n• Best for: Single user", MID_BLUE),
    ("Option B\nDocker Deploy", "• Build Docker image\n• docker run -p 8501:8501\n• Portable across environments\n• Best for: Consistent installs", DARK_BLUE),
    ("Option C\nNetwork Share", "• --server.address=0.0.0.0\n• Accessible on intranet\n• Firewall-restricted access\n• Best for: Team access", RGBColor(0, 80, 160)),
]
for i, (title, desc, color) in enumerate(options):
    left = 0.5 + i * 4.2
    add_rect(s, left, 1.5, 3.8, 4.0, color)
    add_textbox(s, title, left + 0.15, 1.6,  3.5, 0.9, font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(s, desc,  left + 0.15, 2.5,  3.5, 2.8, font_size=12, color=WHITE)

add_textbox(s, "All options run on intranet only — no customer data exposed to external systems.",
            0.4, 5.7, 12.5, 0.5, font_size=12, italic=True, color=MID_BLUE, align=PP_ALIGN.CENTER)


# ── SLIDE 11 — Business Impact ────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 1.2, DARK_BLUE)
add_textbox(s, "Business Impact", 0.4, 0.2, 12, 0.8, font_size=28, bold=True, color=WHITE)

impacts = [
    ("⏱️", "Faster Policy Lookup",    "Reduces search time from 5–10 min to under 30 sec"),
    ("✅", "Consistent Answers",      "Every employee gets the same approved policy content"),
    ("📉", "Reduced Compliance Risk", "Answers sourced from official, reviewed policy documents"),
    ("🚀", "Employee Onboarding",     "New employees self-serve policy questions from day 1"),
    ("🔒", "Zero Data Risk",          "No customer data, no external APIs, no PII exposure"),
    ("💰", "Low Cost Deployment",     "No cloud costs — runs on any internal server"),
]
for i, (icon, title, detail) in enumerate(impacts):
    left = 0.4 + (i % 2) * 6.4
    top  = 1.4 + (i // 2) * 1.7
    add_rect(s, left, top, 6.0, 1.5, LIGHT_BLUE)
    add_textbox(s, f"{icon}  {title}", left + 0.15, top + 0.1,  5.7, 0.55, font_size=14, bold=True, color=DARK_BLUE)
    add_textbox(s, detail,             left + 0.15, top + 0.65, 5.7, 0.6,  font_size=12, color=DARK_GRAY)


# ── SLIDE 12 — Thank You / Q&A ────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BLUE)
add_rect(s, 0, 3.0, 13.33, 0.06, GOLD)

add_textbox(s, "🏦", 5.9, 0.5, 1.5, 1.2, font_size=48, align=PP_ALIGN.CENTER, color=WHITE)
add_textbox(s, "Thank You", 0.5, 1.6, 12.3, 1.2,
            font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, "Internal Bank Employee Assistant — Capstone Project",
            0.5, 2.5, 12.3, 0.6, font_size=16, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_textbox(s, "Questions & Discussion",
            0.5, 3.4, 12.3, 0.8, font_size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_textbox(s, "Source Code  •  Architecture Diagram  •  All Documents available in deliverables/",
            0.5, 4.5, 12.3, 0.6, font_size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save("/home/labuser/Day10/deliverables/05_Presentation_Deck.pptx")
print("✅ Presentation Deck (.pptx) created.")
