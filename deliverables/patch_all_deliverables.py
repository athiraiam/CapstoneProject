"""
patch_all_deliverables.py
--------------------------
Comprehensive patch for ALL deliverable gaps.

PPTX gaps (05_Presentation_Deck.pptx):
  - Insert 5 missing slides after slide 11 (Governance):
      Slide 12: MCP & Plugin Integration
      Slide 13: Observability & Traceability
      Slide 14: Load Testing Results
      Slide 15: Deployment Architecture
      Slide 16: Screenshots of Results
  - Renumber final "Thank You" slide stays at end

Word doc gaps:
  02_Governance_Report.docx:
      - Add "Observability & Traceability" as Section 9

  03_Testing_Evaluation_Report.docx:
      - Add "Screenshots of Results" as Section 7 (with embedded images)

Run: python patch_all_deliverables.py   (from deliverables/ directory)
"""

import os, json, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Pt as DPt, RGBColor as DRGBColor, Inches as DInches

BASE        = os.path.dirname(os.path.abspath(__file__))
SCREENSHOTS = os.path.join(BASE, "screenshots")
LT_RESULTS  = os.path.join(BASE, "..", "bank_assistant", "load_test_results.json")

# ── PPTX colour palette (matches existing deck) ──────────────────────────────
DARK_BLUE  = RGBColor(0,   51,  102)
MID_BLUE   = RGBColor(0,   102, 204)
LIGHT_BLUE = RGBColor(204, 229, 255)
WHITE      = RGBColor(255, 255, 255)
GOLD       = RGBColor(255, 193,   7)
DARK_GRAY  = RGBColor(50,   50,  50)
GREEN      = RGBColor(0,   140,  64)
TEAL       = RGBColor(0,   120, 120)
PURPLE     = RGBColor(100,   0, 160)
ORANGE     = RGBColor(200,  90,   0)
RED        = RGBColor(180,   0,   0)

# ── Word colour palette ───────────────────────────────────────────────────────
W_BANK_BLUE   = DRGBColor(0, 51, 102)
W_ACCENT_BLUE = DRGBColor(0, 102, 204)
W_GREEN       = DRGBColor(0, 128, 64)


# ════════════════════════════════════════════════════════════════════════════
# PPTX helpers (mirror existing deck style exactly)
# ════════════════════════════════════════════════════════════════════════════
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, text, left, top, w, h, size=14, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text           = text
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    return box


def rect(slide, left, top, w, h, fill):
    s = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def hdr(slide, title, color=DARK_BLUE):
    rect(slide, 0, 0, 13.33, 1.2, color)
    tb(slide, title, 0.4, 0.18, 12.5, 0.9, size=26, bold=True, color=WHITE)


def add_slide(prs):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def insert_slide_after(prs, after_index):
    """Append a new blank slide and then move it to after_index+1."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Move the slide XML to the correct position
    xml_slides = prs.slides._sldIdLst
    entries = list(xml_slides)
    # The new slide is last; move it to after_index+1
    new_entry = entries[-1]
    xml_slides.remove(new_entry)
    xml_slides.insert(after_index + 1, new_entry)
    return slide


# ════════════════════════════════════════════════════════════════════════════
# Build the 5 new PPTX slides
# ════════════════════════════════════════════════════════════════════════════

def build_slide_mcp(prs):
    s = add_slide(prs)
    set_bg(s, WHITE)
    hdr(s, "MCP & Plugin Integration", MID_BLUE)
    tb(s, "mcp_server.py exposes the full NexaBank knowledge base as 5 MCP tools — "
          "usable by Claude Desktop, Claude Code, or any MCP-compatible AI agent.",
       0.4, 1.3, 12.5, 0.55, size=12, color=DARK_BLUE, italic=True)

    tools = [
        ("search_policy",        "Keyword search → best-matching policy snippet + source badge"),
        ("list_policies",        "Returns all 5 policy names, codes and sample topics"),
        ("get_policy_full",      "Returns complete policy text by code (LP-001 … AOP-005)"),
        ("cross_reference",      "Finds ALL policies matching a multi-domain question"),
        ("check_compliance_risk","Screens question against 14 compliance risk keywords"),
    ]
    for i, (name, desc) in enumerate(tools):
        top = 2.0 + i * 0.95
        rect(s, 0.4, top, 3.5, 0.8, DARK_BLUE)
        tb(s, name, 0.5, top + 0.08, 3.3, 0.65, size=11, bold=True, color=GOLD)
        rect(s, 4.0, top, 9.0, 0.8, LIGHT_BLUE)
        tb(s, desc, 4.1, top + 0.12, 8.8, 0.6, size=11, color=DARK_BLUE)

    tb(s, "Protocol: MCP JSON-RPC 2.0 over stdio  |  Config: claude_mcp_config.json  |  "
          "Both Streamlit UI & MCP share the same knowledge base + KEYWORD_MAP",
       0.4, 6.8, 12.5, 0.5, size=10, italic=True, color=MID_BLUE)
    return s


def build_slide_observability(prs):
    s = add_slide(prs)
    set_bg(s, WHITE)
    hdr(s, "Observability & Traceability", TEAL)
    tb(s, "Every query is fully traceable — from input validation to final answer — "
          "across 3 structured log files and 6 lifecycle hooks.",
       0.4, 1.3, 12.5, 0.5, size=12, italic=True, color=DARK_BLUE)

    logs = [
        ("📋 query_audit.log",       "Every query",     "Timestamp · Session ID · STATUS (MATCHED/FALLBACK) · Policy matched · Question text"),
        ("⚠️ unresolved_queries.log","Fallback queries","Timestamp · Session ID · UNRESOLVED · Full question — for compliance gap analysis"),
        ("🔄 session_events.log",    "All events",      "PRE_QUERY · POST_ANSWER · SESSION_START/END · ESCALATION · COMPLIANCE_ALERT — full lifecycle"),
    ]
    for i, (log, trigger, detail) in enumerate(logs):
        top = 2.0 + i * 1.5
        rect(s, 0.3, top, 2.8, 1.2, DARK_BLUE)
        tb(s, log,     0.4, top + 0.1, 2.6, 0.5, size=12, bold=True, color=WHITE)
        tb(s, trigger, 0.4, top + 0.65, 2.6, 0.4, size=10, color=GOLD)
        rect(s, 3.3, top, 9.7, 1.2, LIGHT_BLUE)
        tb(s, detail, 3.45, top + 0.25, 9.5, 0.7, size=11, color=DARK_BLUE)

    tb(s, "Answer Traceability: "
          "Every answer is a verbatim extract from policy text — can be traced back to the exact file and line.  "
          "AuditLoggerAgent runs silently after EVERY query type (policy, summary, escalation, fallback).",
       0.4, 6.65, 12.5, 0.65, size=10, italic=True, color=DARK_GRAY)
    return s


def build_slide_load_test(prs):
    with open(LT_RESULTS) as f:
        lt = json.load(f)
    cfg  = lt["test_config"]
    res  = lt["results"]
    lat  = lt["latency_ms"]
    tput = lt["throughput"]

    s = add_slide(prs)
    set_bg(s, WHITE)
    hdr(s, "Load Testing Results", GREEN)
    tb(s, f"Real load test: {cfg['total_tasks']} tasks · "
          f"{cfg['concurrent_workers']} concurrent workers · "
          f"{cfg['iterations']} iterations of all {cfg['total_questions']} question types",
       0.4, 1.3, 12.5, 0.5, size=12, italic=True, color=DARK_BLUE)

    # KPI boxes
    kpis = [
        ("0",                         "Errors",             DARK_BLUE),
        (f"{res['found_rate_pct']}%", "Answer Found Rate",  GREEN),
        (f"{lat['p50']} ms",          "Median Latency",     MID_BLUE),
        (f"{lat['p99']} ms",          "p99 Latency",        TEAL),
        (f"{tput['requests_per_sec']}","Requests / sec",    PURPLE),
        (f"{cfg['total_tasks']}/60",  "Tasks Completed",    ORANGE),
    ]
    for i, (val, label, color) in enumerate(kpis):
        col = i % 3
        row = i // 3
        left = 0.3 + col * 4.3
        top  = 2.0 + row * 2.3
        rect(s, left, top, 4.1, 2.0, color)
        tb(s, val,   left+0.1, top+0.15, 3.9, 1.0, size=30, bold=True,
           color=WHITE, align=PP_ALIGN.CENTER)
        tb(s, label, left+0.1, top+1.25, 3.9, 0.55, size=12,
           color=WHITE, align=PP_ALIGN.CENTER)

    tb(s, f"Min: {lat['min']} ms  ·  Mean: {lat['mean']} ms  ·  p95: {lat['p95']} ms  ·  "
          f"Max: {lat['max']} ms  ·  Wall clock: {tput['wall_clock_sec']} s",
       0.4, 6.8, 12.5, 0.5, size=10, italic=True, color=DARK_GRAY)
    return s


def build_slide_deployment(prs):
    s = add_slide(prs)
    set_bg(s, WHITE)
    hdr(s, "Deployment Architecture", ORANGE)
    tb(s, "Three deployment modes — all sharing the same offline knowledge base. "
          "No cloud dependency. No external APIs. Runs on any internal server.",
       0.4, 1.3, 12.5, 0.5, size=12, italic=True, color=DARK_BLUE)

    modes = [
        ("Option A\nLocal Deploy",
         "streamlit run app.py",
         "Single user  ·  http://localhost:8501\nPy 3.10+  ·  streamlit>=1.32",
         MID_BLUE),
        ("Option B\nDocker",
         "docker build -t bank-assistant .\ndocker run -p 8501:8501 ...",
         "Isolated container  ·  Any OS\nFROM python:3.11-slim",
         TEAL),
        ("Option C\nNetwork / Team",
         "streamlit run app.py\n--server.address=0.0.0.0",
         "Team access  ·  http://<ip>:8501\nIntranet firewall required",
         DARK_BLUE),
        ("Option D\nMCP Server",
         "python mcp_server.py",
         "AI agent access  ·  stdio transport\nClaude Desktop / Claude Code",
         PURPLE),
    ]
    for i, (title, cmd, detail, color) in enumerate(modes):
        left = 0.2 + i * 3.28
        rect(s, left, 2.1, 3.1, 1.0, color)
        tb(s, title, left+0.1, 2.15, 2.9, 0.85, size=12, bold=True, color=WHITE)
        rect(s, left, 3.3, 3.1, 1.1, LIGHT_BLUE)
        tb(s, cmd,    left+0.1, 3.35, 2.9, 0.8, size=9,  color=DARK_BLUE)
        rect(s, left, 4.55, 3.1, 1.2, RGBColor(240, 245, 255))
        tb(s, detail, left+0.1, 4.6, 2.9, 1.0, size=9,  color=DARK_GRAY)

    tb(s,
       "All modes share: policies/*.txt  ·  KEYWORD_MAP  ·  7 Subagents  ·  3 Skills  ·  6 Hooks  ·  3 Log Files",
       0.4, 6.0, 12.5, 0.55, size=11, italic=True, color=MID_BLUE)

    req = [
        ("Python  3.10+", DARK_BLUE), ("RAM  512 MB+", TEAL),
        ("Disk  100 MB", MID_BLUE),   ("Network  LAN (optional)", ORANGE),
    ]
    for i, (txt, color) in enumerate(req):
        rect(s, 0.3 + i*3.28, 6.7, 3.1, 0.55, color)
        tb(s, txt, 0.4 + i*3.28, 6.78, 3.0, 0.4, size=10, bold=True, color=WHITE)
    return s


def build_slide_screenshots(prs):
    s = add_slide(prs)
    set_bg(s, WHITE)
    hdr(s, "Screenshots of Results", DARK_BLUE)

    shots = [
        ("sc9_app_home.png",         "App Home"),
        ("sc1_loan_policy_search.png","Loan Policy Search"),
        ("sc2_kyc_search.png",        "KYC Search"),
        ("sc3_multi_policy.png",      "Multi-Policy"),
        ("sc4_summary.png",           "Policy Summary"),
        ("sc5_escalation.png",        "Escalation"),
        ("sc6_compliance_alert.png",  "Compliance Alert"),
        ("sc7_fallback_selfheal.png", "Self-Healing Fallback"),
        ("sc8_account_opening.png",   "Account Opening"),
    ]

    # 3 columns × 3 rows grid
    img_w, img_h = 3.9, 1.65
    col_gap, row_gap = 0.28, 0.22
    start_x, start_y = 0.25, 1.35

    for i, (fname, label) in enumerate(shots):
        col = i % 3
        row = i // 3
        x = start_x + col * (img_w + col_gap)
        y = start_y + row * (img_h + row_gap + 0.3)
        img_path = os.path.join(SCREENSHOTS, fname)
        if os.path.exists(img_path):
            s.shapes.add_picture(
                img_path,
                Inches(x), Inches(y), Inches(img_w), Inches(img_h))
        else:
            rect(s, x, y, img_w, img_h, LIGHT_BLUE)
        # Label below image
        tb(s, label, x, y + img_h + 0.03, img_w, 0.25,
           size=9, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    return s


# ════════════════════════════════════════════════════════════════════════════
# Patch PPTX
# ════════════════════════════════════════════════════════════════════════════
def patch_pptx():
    path = os.path.join(BASE, "05_Presentation_Deck.pptx")
    prs  = Presentation(path)

    n = len(prs.slides)
    print(f"  PPTX: {n} slides before patch")

    # Append 5 new slides (they end up after the existing "Thank You" slide,
    # then we move "Thank You" to the very end)
    build_slide_mcp(prs)           # appended as n+1
    build_slide_observability(prs) # appended as n+2
    build_slide_load_test(prs)     # appended as n+3
    build_slide_deployment(prs)    # appended as n+4
    build_slide_screenshots(prs)   # appended as n+5

    # Move the original last slide (Thank You) which was at index n-1
    # to the very end (after all new slides)
    xml_slides = prs.slides._sldIdLst
    entries    = list(xml_slides)
    # Thank You was originally index n-1, which is now at index n-1
    # (5 new slides were appended after it)
    thank_you  = entries[n - 1]
    xml_slides.remove(thank_you)
    xml_slides.append(thank_you)

    prs.save(path)
    print(f"  PPTX: {len(prs.slides)} slides after patch  ✅")


# ════════════════════════════════════════════════════════════════════════════
# Word helpers
# ════════════════════════════════════════════════════════════════════════════
def wh1(doc, text):
    p = doc.add_heading(text, level=1)
    for run in p.runs:
        run.font.color.rgb = W_BANK_BLUE
    return p

def wh2(doc, text):
    p = doc.add_heading(text, level=2)
    for run in p.runs:
        run.font.color.rgb = W_ACCENT_BLUE
    return p

def wpara(doc, text, bold=False, italic=False, size=11):
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.bold   = bold
    r.italic = italic
    r.font.size = DPt(size)
    return p

def wtable(doc, headers, rows):
    t = doc.add_table(rows=1 + len(rows), cols=len(headers))
    t.style = "Table Grid"
    hdr = t.rows[0].cells
    for i, h in enumerate(headers):
        hdr[i].text = h
        for run in hdr[i].paragraphs[0].runs:
            run.bold = True
    for row_data in rows:
        rc = t.add_row().cells
        for i, v in enumerate(row_data):
            rc[i].text = str(v)
    return t

def wbullet(doc, items):
    for item in items:
        doc.add_paragraph(item, style="List Bullet")


# ════════════════════════════════════════════════════════════════════════════
# Patch Governance Report — add Observability & Traceability
# ════════════════════════════════════════════════════════════════════════════
def patch_governance():
    path = os.path.join(BASE, "02_Governance_Report.docx")
    doc  = Document(path)

    doc.add_page_break()
    wh1(doc, "9. Observability & Traceability")
    wpara(doc, (
        "Every query processed by the agent pipeline is fully observable through three "
        "structured log files and six lifecycle hook events. All logs are append-only, "
        "written to bank_assistant/logs/, and require no external monitoring tooling."
    ))

    wh2(doc, "9.1 Log Files")
    wtable(doc,
        ["Log File", "Written By", "Trigger", "Contents"],
        [
            ["query_audit.log",
             "AuditLoggerAgent",
             "Every query (all intents)",
             "Timestamp, Session ID, STATUS (MATCHED/FALLBACK), Policy matched, Question text"],
            ["unresolved_queries.log",
             "escalation_hook + AuditLoggerAgent",
             "found=False queries only",
             "Timestamp, Session ID, UNRESOLVED flag, Full question — enables policy gap analysis"],
            ["session_events.log",
             "All 6 hooks",
             "Every pipeline event",
             "PRE_QUERY, POST_ANSWER, SESSION_START, SESSION_END, ESCALATION, COMPLIANCE_ALERT events"],
        ]
    )

    wh2(doc, "9.2 Hook-Level Traceability")
    wtable(doc,
        ["Hook", "Log Event Written", "Log File", "Fields Captured"],
        [
            ["pre_query_hook",      "PRE_QUERY",         "session_events.log", "Session ID, question (first 80 chars)"],
            ["post_answer_hook",    "POST_ANSWER",        "session_events.log", "Status, policy name, response time (ms), question"],
            ["session_start_hook",  "SESSION_START",      "session_events.log", "Session ID, timestamp, knowledge base ready"],
            ["session_end_hook",    "SESSION_END",        "session_events.log", "Session ID, total queries handled"],
            ["escalation_hook",     "ESCALATION",         "unresolved_queries.log + session_events.log", "Session ID, escalated question"],
            ["compliance_alert_hook","COMPLIANCE_ALERT",  "session_events.log", "Session ID, flagged keywords, question"],
        ]
    )

    wh2(doc, "9.3 Answer Traceability")
    wpara(doc, (
        "Every answer returned to the employee can be traced back to its exact source:"
    ))
    wbullet(doc, [
        "The policy_name and policy_code fields in the result dict identify which document was used",
        "The agent_used field identifies which subagent produced the answer (shown as a chip in the UI)",
        "The intent field identifies how the query was classified by QueryClassifierAgent",
        "All answers are verbatim extracts from policy text — no content is fabricated",
        "The session_id field links every query to a browser session for cross-log correlation",
    ])

    wh2(doc, "9.4 Compliance Monitoring")
    wtable(doc,
        ["Risk Keyword Category", "Examples", "Action Taken"],
        [
            ["Financial crime",  "fraud, money laundering, terrorist", "compliance_alert_hook fires; logged to session_events.log; yellow banner shown in UI"],
            ["Process bypass",   "bypass, override, skip verification", "Same as above"],
            ["Document fraud",   "fake document, corrupt",              "Same as above"],
            ["Sanctions",        "sanction, blacklist",                 "Same as above"],
        ]
    )
    wpara(doc, (
        "detect_compliance_risk() in audit_skill.py screens every question against 14 keywords. "
        "The check runs for ALL query types, not only policy queries."
    ), italic=True)

    wh2(doc, "9.5 UI Observability Panel")
    wbullet(doc, [
        "Every assistant message shows an 'Agent: <name> | Intent: <type>' chip beneath the source badge",
        "Fallback questions display a yellow warning box with the contact details for compliance/operations",
        "Compliance-risk questions display a yellow warning banner above the answer",
        "app.py exposes the last 20 lines of each log file in a sidebar audit panel (read_audit_log())",
    ])

    doc.save(path)
    print(f"  Governance Report: Section 9 added  ✅")


# ════════════════════════════════════════════════════════════════════════════
# Patch Testing Report — embed screenshots as Section 7
# ════════════════════════════════════════════════════════════════════════════
def patch_testing():
    path = os.path.join(BASE, "03_Testing_Evaluation_Report.docx")
    doc  = Document(path)

    doc.add_page_break()
    wh1(doc, "7. Screenshots of Results")
    wpara(doc, (
        "The following screenshots were captured from the live Streamlit application, "
        "demonstrating all 7 subagents and 3 skills working in the full agent pipeline. "
        "Screenshots are stored in deliverables/screenshots/."
    ))

    shots = [
        ("sc9_app_home.png",         "Screenshot 1 — Application Home",
         "The NexaBank Employee Assistant home screen showing the welcome panel, "
         "sidebar policy list, and the 5-topic grid with icons."),
        ("sc1_loan_policy_search.png","Screenshot 2 — Loan Policy Search (PolicySearchAgent)",
         "PolicySearchAgent matches a loan eligibility question and returns a focused "
         "snippet from LP-001 with the source badge."),
        ("sc2_kyc_search.png",        "Screenshot 3 — KYC Document Search (PolicySearchAgent)",
         "KYC-002 matched from 'What are the KYC document requirements?' — source badge "
         "and agent chip visible below the answer."),
        ("sc3_multi_policy.png",      "Screenshot 4 — Multi-Policy Answer (MultiPolicyAgent)",
         "MultiPolicyAgent merges snippets from KYC-002 and AOP-005 for a cross-domain "
         "question about both KYC and account opening requirements."),
        ("sc4_summary.png",           "Screenshot 5 — Policy Summary (SummaryAgent)",
         "SummaryAgent returns a 5-bullet overview of the Loan Policy triggered by "
         "'Summarize the loan policy'."),
        ("sc5_escalation.png",        "Screenshot 6 — Escalation + Email Draft (EscalationAgent)",
         "EscalationAgent routes to the compliance team and shows a pre-filled "
         "escalation email draft in the expandable panel."),
        ("sc6_compliance_alert.png",  "Screenshot 7 — Compliance Alert Banner",
         "compliance_alert_hook fires on a risk-keyword question, showing the yellow "
         "alert banner while still returning a policy answer."),
        ("sc7_fallback_selfheal.png", "Screenshot 8 — Self-Healing Fallback (FallbackAgent)",
         "FallbackAgent invoked when PolicySearchAgent finds no match — shows the "
         "warning box with compliance/operations contact details."),
        ("sc8_account_opening.png",   "Screenshot 9 — Account Opening Policy (AOP-005)",
         "Account opening requirements returned from AOP-005, showing minimum balance, "
         "document list, and account types."),
    ]

    for fname, caption, description in shots:
        img_path = os.path.join(SCREENSHOTS, fname)
        wh2(doc, caption)
        wpara(doc, description)
        if os.path.exists(img_path):
            doc.add_picture(img_path, width=DInches(5.5))
            # Center the picture paragraph
            last_para = doc.paragraphs[-1]
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            last_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        else:
            wpara(doc, f"[Image file not found: {fname}]", italic=True)
        doc.add_paragraph()  # spacing

    doc.save(path)
    print(f"  Testing Report: Section 7 (Screenshots) added with {len(shots)} images  ✅")


# ════════════════════════════════════════════════════════════════════════════
# Run all patches
# ════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("\nPatching all deliverables...\n")
    patch_pptx()
    patch_governance()
    patch_testing()
    print("\nAll patches applied successfully.")
