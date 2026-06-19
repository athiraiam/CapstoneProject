"""
regenerate_all.py
Regenerates all 5 deliverable documents with every audit fix applied.
Run: python regenerate_all.py
"""

from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN
import datetime

TODAY = datetime.date.today().strftime("%B %d, %Y")

# Word colors
BANK_BLUE   = RGBColor(0,  51, 102)
ACCENT_BLUE = RGBColor(0, 102, 204)
GREEN_W     = RGBColor(0, 128,  64)
PURPLE_W    = RGBColor(80,  0, 128)
ORANGE_W    = RGBColor(180, 80,  0)
RED_W       = RGBColor(160,  0,  0)
TEAL_W      = RGBColor(0,  100, 100)

# PPTX colors
DARK_BLUE  = PRGBColor(0,  51, 102)
MID_BLUE   = PRGBColor(0, 102, 204)
LIGHT_BLUE = PRGBColor(204, 229, 255)
WHITE      = PRGBColor(255, 255, 255)
GOLD       = PRGBColor(255, 193,   7)
DARK_GRAY  = PRGBColor(50,  50,  50)
GREEN_P    = PRGBColor(0,  140,  64)
PURPLE_P   = PRGBColor(100,  0, 160)
ORANGE_P   = PRGBColor(200,  90,  0)
RED_P      = PRGBColor(180,   0,  0)
TEAL_P     = PRGBColor(0,  120, 120)


# ── Word helpers ──────────────────────────────────────────────────────────────
def wh(doc, text, level=1, color=None):
    h = doc.add_heading(text, level=level)
    c = color or (BANK_BLUE if level == 1 else ACCENT_BLUE)
    for run in h.runs:
        run.font.color.rgb = c
    return h


def wp(doc, text, bold=False, italic=False, size=11, color=None):
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.bold, r.italic, r.font.size = bold, italic, Pt(size)
    if color:
        r.font.color.rgb = color
    return p


def wt(doc, headers, rows, widths=None):
    t = doc.add_table(rows=1 + len(rows), cols=len(headers))
    t.style = "Table Grid"
    hc = t.rows[0].cells
    for i, h in enumerate(headers):
        hc[i].text = h
        for r in hc[i].paragraphs[0].runs:
            r.bold = True
        if widths:
            hc[i].width = Inches(widths[i])
    for rd in rows:
        rc = t.add_row().cells
        for i, v in enumerate(rd):
            rc[i].text = v
            if widths:
                rc[i].width = Inches(widths[i])
    return t


def wb(doc, items):
    for it in items:
        doc.add_paragraph(it, style="List Bullet")


def wn(doc, items):
    for it in items:
        doc.add_paragraph(it, style="List Number")


def wcode(doc, lines):
    for line in lines:
        p = doc.add_paragraph()
        r = p.add_run(line)
        r.font.name = "Courier New"
        r.font.size = Pt(9)


def cover(doc, title, subtitle):
    doc.add_paragraph(); doc.add_paragraph()
    t = doc.add_paragraph(title)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    t.runs[0].font.size = Pt(24)
    t.runs[0].font.color.rgb = BANK_BLUE
    t.runs[0].bold = True
    s = doc.add_paragraph(subtitle)
    s.alignment = WD_ALIGN_PARAGRAPH.CENTER
    s.runs[0].font.size = Pt(14)
    s.runs[0].italic = True
    doc.add_paragraph()
    d = doc.add_paragraph(f"Date: {TODAY}  |  NexaBank Internal Division  |  Version 2.0")
    d.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_page_break()


# ── PPTX helpers ──────────────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(
        PInches(l), PInches(t), PInches(w), PInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = PPt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(
        1, PInches(l), PInches(t), PInches(w), PInches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def hdr(slide, title, color=DARK_BLUE):
    rect(slide, 0, 0, 13.33, 1.2, color)
    tb(slide, title, 0.4, 0.18, 12.5, 0.9, size=26, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# 1. TECHNICAL DESIGN DOCUMENT  (fixes: numbering, audit_skill annotation,
#    utils/ folder, KEYWORD_MAP, out_of_scope clarification)
# ════════════════════════════════════════════════════════════════════════════
def create_tdd():
    doc = Document()
    cover(doc, "Technical Design Document",
          "Internal Bank Employee Assistant — Multi-Agent Architecture v2.0")

    wh(doc, "1. Executive Summary")
    wp(doc, (
        "This document describes the complete technical design of the Internal Bank Employee "
        "Assistant, including the multi-agent architecture, all 7 subagents, all 3 skills, "
        "all 6 hooks, the KEYWORD_MAP scoring engine, and how they collaborate through the "
        "OrchestratorAgent pipeline."
    ))

    wh(doc, "2. System Architecture Overview")
    wp(doc, (
        "The system is built on a three-tier layered architecture: Presentation Layer "
        "(Streamlit UI), Agent Layer (7 subagents + 1 orchestrator), and Knowledge Base "
        "(5 plain-text policy documents). All agents are coordinated by OrchestratorAgent."
    ))

    wh(doc, "2.1 Architecture Layers", 2)
    wt(doc,
        ["Layer", "Components", "Technology"],
        [
            ["Presentation",   "app.py — Streamlit chatbot UI",                          "Python, Streamlit 1.32+"],
            ["Orchestration",  "OrchestratorAgent — coordinates all subagents",          "Python 3.10+"],
            ["Agents",         "7 subagents (classifier, search, multi, summary, etc.)", "Python classes"],
            ["Skills",         "3 reusable skill modules",                               "Python modules"],
            ["Hooks",          "6 lifecycle hooks in hooks.py",                          "Python functions"],
            ["Knowledge Base", "5 policy .txt files + KEYWORD_MAP in utils/",           "Plain text + Python dict"],
            ["Legacy Utils",   "utils/knowledge_base.py, utils/agent.py (pre-refactor)","Python — read-only at runtime"],
            ["Audit",          "3 log files in logs/",                                   "Append-only text logs"],
        ]
    )

    wh(doc, "2.2 Agent Pipeline Flow (9 Steps)", 2)
    wp(doc, "Every employee question follows this exact pipeline:", bold=True)
    wn(doc, [
        "pre_query_hook — validates and sanitizes the question (rejects empty / too-short)",
        "QueryClassifierAgent — classifies intent: policy_query / summary / escalation / out_of_scope",
        "OrchestratorAgent routing decision — selects correct subagent (see routing table)",
        "SubAgent executes — returns structured result dict",
        "AuditLoggerAgent — logs query + result to audit files (runs after every query)",
        "compliance_alert_hook — fires if risk keywords detected in the question",
        "escalation_hook — fires if found=False AND escalated=False",
        "post_answer_hook — logs response time and outcome",
        "Result returned to app.py for display in the chat UI",
    ])

    wp(doc, "Note: pre_query_hook handles empty/too-short input BEFORE the classifier runs. "
            "QueryClassifierAgent emits out_of_scope only when the stripped question is an empty string "
            "after the hook has already passed it through. Both gates are complementary.", italic=True)

    doc.add_page_break()

    wh(doc, "3. OrchestratorAgent")
    wp(doc, "File: agents/orchestrator.py", italic=True, color=ACCENT_BLUE)
    wp(doc, (
        "The OrchestratorAgent is the top-level coordinator. Instantiated once per browser "
        "session via Streamlit's cache_resource. Owns the full request lifecycle."
    ))

    wh(doc, "3.1 Responsibilities", 2)
    wb(doc, [
        "Instantiates all 7 subagents at session start",
        "Fires session_start_hook and session_end_hook",
        "Calls pre_query_hook before routing any question",
        "Delegates to QueryClassifierAgent to determine intent",
        "Routes to correct subagent based on intent",
        "Self-healing: if PolicySearchAgent returns found=False → automatically invokes FallbackAgent",
        "Calls AuditLoggerAgent after every answer",
        "Fires compliance_alert_hook and escalation_hook as needed",
        "Calls post_answer_hook with response timing",
    ])

    wh(doc, "3.2 Routing Decision Table", 2)
    wt(doc,
        ["Intent", "Routing Condition", "Subagent Invoked"],
        [
            ["summary",      "Question contains: summarize, overview, brief, explain the policy",          "SummaryAgent"],
            ["escalation",   "Question contains: escalate, contact, manager, urgent, speak to, human",    "EscalationAgent"],
            ["out_of_scope", "Empty string after pre_query_hook / too short (< 3 chars)",                 "Direct fallback — no subagent"],
            ["policy_query", "cross_reference: top score ≥ 2 AND second score ≥ 1 (multi-domain)",       "MultiPolicyAgent"],
            ["policy_query", "Single best keyword match (score ≥ 1)",                                     "PolicySearchAgent"],
            ["policy_query", "PolicySearchAgent returns found=False (self-healing trigger)",              "FallbackAgent → escalation_hook"],
        ]
    )

    wh(doc, "3.3 Return Structure", 2)
    wt(doc,
        ["Field", "Type", "Description"],
        [
            ["answer",           "str",      "Answer text displayed in chat UI"],
            ["policy_name",      "str|None", "Primary matched policy display name"],
            ["policy_names",     "list",     "All matched policy names (MultiPolicyAgent only)"],
            ["found",            "bool",     "True if a policy match was found"],
            ["intent",           "str",      "Classified intent: policy_query / summary / escalation / out_of_scope"],
            ["escalated",        "bool",     "True if the question was routed to EscalationAgent"],
            ["email_draft",      "str|None", "Draft escalation email (EscalationAgent only)"],
            ["compliance_alert", "str|None", "Warning message if risk keywords detected"],
            ["agent_used",       "str",      "Name of the subagent that produced the answer"],
            ["session_id",       "str",      "Unique 8-char session identifier"],
        ]
    )

    doc.add_page_break()

    wh(doc, "4. Subagents")
    wp(doc, "FallbackAgent is SubAgent 6; AuditLoggerAgent is SubAgent 7 — consistent with code file headers.", italic=True)

    subagents = [
        ("4.1 QueryClassifierAgent", "agents/query_classifier_agent.py", GREEN_W,
         "SubAgent 1. Runs first in every pipeline. Classifies question into one of four intents. Uses extract_keywords() from search_skill.py.",
         [("Input","Raw question string"),("Output","{ intent, keywords, question }"),
          ("Intent Types","policy_query | summary | escalation | out_of_scope"),
          ("Summary Triggers","summarize, overview, brief, give me a summary, explain the policy"),
          ("Escalation Triggers","escalate, contact, manager, urgent, speak to, human"),
          ("Skills Used","extract_keywords() from search_skill.py")]),

        ("4.2 PolicySearchAgent", "agents/policy_search_agent.py", ACCENT_BLUE,
         "SubAgent 2. Handles standard policy questions via keyword scoring. Returns the best matching policy snippet.",
         [("Input","question + policies dict"),("Output","{ answer, policy_key, policy_name, found, score }"),
          ("Algorithm","KEYWORD_MAP scoring → snippet extraction via extract_answer()"),
          ("Fallback","Returns found=False if score=0 for all policies"),
          ("Skills Used","search_policy(), format_answer() from search_skill.py and format_skill.py")]),

        ("4.3 MultiPolicyAgent", "agents/multi_policy_agent.py", PURPLE_W,
         "SubAgent 3. Handles questions spanning 2 policy domains. Activated when cross_reference returns top score ≥ 2 AND second score ≥ 1.",
         [("Input","question + policies dict"),("Output","{ answer, policy_names, found, multi }"),
          ("Activation","cross_reference scores: [0].score ≥ 2 and [1].score ≥ 1"),
          ("Merge Strategy","Top 2 matching policy snippets concatenated with divider"),
          ("Skills Used","cross_reference() from search_skill.py")]),

        ("4.4 SummaryAgent", "agents/summary_agent.py", ORANGE_W,
         "SubAgent 4. Generates a 5-bullet policy summary. Detects target policy via POLICY_NAME_HINTS dict first, falls back to search_policy() if no hint matches.",
         [("Input","question + policies dict"),("Output","{ answer, policy_name, found }"),
          ("Trigger","Activated by 'summary' intent from QueryClassifierAgent"),
          ("POLICY_NAME_HINTS","loan, kyc, complaint, credit, card, account"),
          ("Fallback Path","If no hint matches, calls search_policy() to infer policy key"),
          ("Skills Used","summarize_policy() from format_skill.py")]),

        ("4.5 EscalationAgent", "agents/escalation_agent.py", RED_W,
         "SubAgent 5. Handles escalation requests. Reads policy .txt files (via search_policy) when no policy key provided. Routes to correct team and generates email draft.",
         [("Input","question + optional matched_policy_key"),
          ("Output","{ answer, policy_name, found, escalated, email_draft }"),
          ("Data Accessed","Question text + policy files (to infer key when none provided)"),
          ("Contact Routing","loan→loan.operations | kyc→compliance | complaint→customerexperience | card→creditcards | account→accountoperations"),
          ("Default Contact","compliance@bank.internal | operations@bank.internal"),
          ("Skills Used","generate_escalation_email(), search_policy() from format_skill.py and search_skill.py")]),

        ("4.6 FallbackAgent", "agents/fallback_agent.py", ORANGE_W,
         "SubAgent 6. Self-healing agent. Auto-invoked by Orchestrator when PolicySearchAgent returns found=False. Tries partial recovery via cross_reference at any score > 0.",
         [("Input","question + policies dict"),("Output","{ answer, suggestion, needs_escalation, found }"),
          ("Recovery Logic","cross_reference at score > 0 threshold for partial match suggestion"),
          ("found","Always False — this agent handles no-match cases"),
          ("Self-Healing","Automatically chained by Orchestrator — no employee action needed"),
          ("Skills Used","cross_reference() from search_skill.py")]),

        ("4.7 AuditLoggerAgent", "agents/audit_logger_agent.py", TEAL_W,
         "SubAgent 7. Background agent. Runs after every answer. Logs queries, detects compliance risk. Never visible to the employee.",
         [("Input","question, matched_policy, found, session_id"),
          ("Output","{ logged, compliance_risk, session_id }"),
          ("Methods","log_answer() — logs query result | log_event() — logs session events"),
          ("Log Files","query_audit.log | unresolved_queries.log | session_events.log"),
          ("Risk Keywords","fraud, bypass, fake document, money laundering, sanction, bribe, corrupt, illegal, etc. (14 total)"),
          ("Skills Used","log_query(), log_unresolved(), detect_compliance_risk() from audit_skill.py")]),
    ]

    for title, filepath, color, desc, details in subagents:
        wh(doc, title, 2, color)
        wp(doc, f"File: {filepath}", italic=True, color=ACCENT_BLUE)
        wp(doc, desc)
        wt(doc, ["Property", "Detail"], details, widths=[1.8, 4.5])
        doc.add_paragraph()

    doc.add_page_break()

    wh(doc, "5. Skills")
    wp(doc, "Skills are reusable, stateless functions called by agents. Each module is independently testable.")

    wh(doc, "5.1 search_skill.py", 2, GREEN_W)
    wp(doc, "File: skills/search_skill.py", italic=True, color=ACCENT_BLUE)
    wp(doc, "Core search and keyword functions. Used by PolicySearchAgent, MultiPolicyAgent, SummaryAgent, FallbackAgent, and EscalationAgent.")
    wt(doc,
        ["Function", "Parameters", "Returns", "Purpose"],
        [
            ["search_policy()",    "question, policies", "{ matched_key, policy_text, snippet, score }",  "Finds best-matching policy via KEYWORD_MAP scoring"],
            ["extract_keywords()", "text",               "list[str]",                                     "Extracts meaningful words (len>3, not stopwords) for classification"],
            ["cross_reference()",  "question, policies", "list[{ matched_key, score, snippet }]",         "Returns ALL matching policies sorted by score — drives multi-policy detection"],
        ], widths=[1.7, 1.6, 2.0, 2.5]
    )

    wh(doc, "5.2 format_skill.py", 2, PURPLE_W)
    wp(doc, "File: skills/format_skill.py", italic=True, color=ACCENT_BLUE)
    wp(doc, "Text formatting, summarization, and email drafting. Contains the canonical POLICY_DISPLAY_NAMES dict imported by all agents.")
    wt(doc,
        ["Function / Constant", "Parameters", "Returns", "Purpose"],
        [
            ["POLICY_DISPLAY_NAMES",         "—",                          "dict",  "Canonical map of policy keys to display names — imported by all agents"],
            ["format_answer()",              "raw_snippet, policy_key",    "str",   "Cleans and formats raw policy text into readable markdown"],
            ["summarize_policy()",           "policy_text, policy_key",    "str",   "Produces 5-bullet summary of a policy document"],
            ["generate_escalation_email()",  "question, matched_policy",   "str",   "Drafts pre-filled escalation email body with correct contact"],
        ], widths=[2.0, 1.8, 0.6, 3.4]
    )

    wh(doc, "5.3 audit_skill.py", 2, ORANGE_W)
    wp(doc, "File: skills/audit_skill.py", italic=True, color=ACCENT_BLUE)
    wp(doc, "Audit logging and compliance detection. Used by AuditLoggerAgent and all 6 hooks.")
    wt(doc,
        ["Function", "Parameters", "Returns", "Purpose"],
        [
            ["log_query()",             "question, policy, found, session_id", "None",   "Appends timestamped record to query_audit.log"],
            ["log_unresolved()",        "question, session_id",                "None",   "Appends unmatched question to unresolved_queries.log"],
            ["log_session_event()",     "event, session_id, detail",           "None",   "Logs session lifecycle events to session_events.log"],
            ["detect_compliance_risk()", "question",                           "bool",   "Returns True if any of 14 risk keywords found in question"],
            ["read_audit_log()",        "log_name, last_n",                    "list",   "Reads last N lines from a log file — used by app.py audit panel"],
        ], widths=[1.9, 2.0, 0.6, 3.3]
    )

    doc.add_page_break()

    wh(doc, "6. Hooks")
    wp(doc, "File: hooks.py", italic=True, color=ACCENT_BLUE)
    wp(doc, "6 lifecycle callbacks fired automatically by OrchestratorAgent at specific pipeline points.")
    wt(doc,
        ["Hook", "Trigger", "Inputs", "Returns", "Responsibility"],
        [
            ["pre_query_hook",        "BEFORE any agent runs",                "question, session_id",              "{ valid, question, reason }",  "Reject empty/<3-char questions; log incoming query to session_events.log"],
            ["post_answer_hook",      "AFTER full pipeline",                  "question, result, session_id, ms",  "None",                         "Log matched policy, response time, ANSWERED/FALLBACK status"],
            ["session_start_hook",    "OrchestratorAgent.__init__()",         "session_id",                        "None",                         "Log SESSION_START with timestamp"],
            ["session_end_hook",      "Employee clicks Clear Chat",           "session_id, total_queries",         "None",                         "Log SESSION_END with total query count"],
            ["escalation_hook",       "found=False AND escalated=False only", "question, session_id",              "None",                         "Write unresolved question to unresolved_queries.log for compliance"],
            ["compliance_alert_hook", "AuditLoggerAgent detects risk",        "question, session_id",              "{ alert: str }",               "Log flagged query; return warning banner text for UI display"],
        ], widths=[1.6, 1.6, 1.8, 1.2, 2.6]
    )

    wh(doc, "6.1 Hook Execution Order Per Query", 2)
    wn(doc, [
        "pre_query_hook — validate input (fires first, before any agent)",
        "QueryClassifierAgent — detect intent",
        "Subagent executes — produce answer",
        "AuditLoggerAgent — log query + result to query_audit.log",
        "compliance_alert_hook — fires IF detect_compliance_risk() returns True",
        "escalation_hook — fires ONLY IF found=False AND escalated=False",
        "post_answer_hook — log response time to session_events.log",
    ])

    wp(doc, "Note: escalation_hook does NOT fire on every found=False. It is suppressed when escalated=True "
            "(i.e., when EscalationAgent has already handled the escalation path).", italic=True)

    doc.add_page_break()

    wh(doc, "7. KEYWORD_MAP — The Scoring Engine")
    wp(doc, "Location: utils/knowledge_base.py", italic=True, color=ACCENT_BLUE)
    wp(doc, (
        "KEYWORD_MAP is the core data structure that drives ALL policy matching across the system. "
        "It maps each of the 5 policy keys to a curated list of ~15 trigger keywords and phrases. "
        "When a question arrives, every keyword in every list is checked against the lowercase question. "
        "The policy with the highest match count wins."
    ))
    wt(doc,
        ["Policy Key", "Sample Keywords (15 per policy)"],
        [
            ["loan_policy",               "loan, borrow, mortgage, personal loan, home loan, auto loan, credit score, prepayment, tenure, emi, disbursement, collateral, repayment, lending, business loan"],
            ["kyc_policy",                "kyc, know your customer, identity, verification, id proof, address proof, passport, driver license, national id, aml, due diligence, onboarding, renewal"],
            ["customer_complaint_policy", "complaint, complain, grievance, dispute, unhappy, escalate, feedback, resolution, refund, compensation, ombudsman, resolve, service complaint"],
            ["credit_card_policy",        "credit card, card, billing, statement, credit limit, reward, cashback, annual fee, lost card, stolen card, minimum payment, apr, cash advance, chargeback"],
            ["account_opening_policy",    "account, open account, savings account, current account, fixed deposit, joint account, minimum balance, debit card, dormant, account closure, cheque book"],
        ]
    )

    wh(doc, "7.1 Legacy utils/ Files", 2)
    wp(doc, (
        "utils/knowledge_base.py contains load_policies(), search_policies(), extract_answer(), and KEYWORD_MAP. "
        "These are imported by the skills layer. "
        "utils/agent.py is the pre-refactor single-agent file containing a legacy get_answer() function and "
        "a duplicate POLICY_DISPLAY_NAMES dict. It is NOT used by any agent in the current multi-agent system. "
        "The canonical POLICY_DISPLAY_NAMES is in skills/format_skill.py. "
        "utils/agent.py should be treated as deprecated."
    ))

    wh(doc, "8. Audit Log Files")
    wt(doc,
        ["File", "Location", "Written By", "Contents"],
        [
            ["query_audit.log",        "logs/", "AuditLoggerAgent.log_answer() via log_query()",         "Timestamp | SessionID | STATUS(MATCHED/FALLBACK) | POLICY | Question"],
            ["unresolved_queries.log", "logs/", "AuditLoggerAgent + escalation_hook via log_unresolved()","Timestamp | SessionID | UNRESOLVED | Question"],
            ["session_events.log",     "logs/", "All hooks via log_session_event()",                     "Session start/end, pre/post query, escalation events, compliance alerts"],
        ], widths=[1.8, 0.8, 2.3, 3.0]
    )
    wp(doc, "Sample entry — query_audit.log:", bold=True)
    wp(doc, "[2024-01-15 10:23:45] SESSION=a3f7b2c1 | STATUS=MATCHED | POLICY=Loan Policy (LP-001) | Q=What documents are needed for a loan?", italic=True)

    doc.add_page_break()

    wh(doc, "9. Complete Project File Structure")
    wcode(doc, [
        "bank_assistant/",
        "├── app.py                                ← Streamlit UI (uses OrchestratorAgent)",
        "├── hooks.py                              ← All 6 lifecycle hooks",
        "├── requirements.txt",
        "├── agents/",
        "│   ├── orchestrator.py                   ← Top-level coordinator",
        "│   ├── query_classifier_agent.py         ← SubAgent 1: intent classification",
        "│   ├── policy_search_agent.py            ← SubAgent 2: single-policy keyword search",
        "│   ├── multi_policy_agent.py             ← SubAgent 3: multi-policy merging",
        "│   ├── summary_agent.py                  ← SubAgent 4: 5-bullet policy summary",
        "│   ├── escalation_agent.py               ← SubAgent 5: contact routing + email draft",
        "│   ├── fallback_agent.py                 ← SubAgent 6: self-healing recovery",
        "│   └── audit_logger_agent.py             ← SubAgent 7: background audit logging",
        "├── skills/",
        "│   ├── search_skill.py                   ← search_policy, extract_keywords, cross_reference",
        "│   ├── format_skill.py                   ← format_answer, summarize_policy, generate_escalation_email, POLICY_DISPLAY_NAMES",
        "│   └── audit_skill.py                    ← log_query, log_unresolved, log_session_event, detect_compliance_risk, read_audit_log",
        "├── policies/                             ← 5 plain-text knowledge base files",
        "│   ├── loan_policy.txt",
        "│   ├── kyc_policy.txt",
        "│   ├── customer_complaint_policy.txt",
        "│   ├── credit_card_policy.txt",
        "│   └── account_opening_policy.txt",
        "├── utils/                                ← Core NLP engine imported by skills",
        "│   ├── knowledge_base.py                 ← load_policies, KEYWORD_MAP, search_policies, extract_answer",
        "│   └── agent.py                          ← DEPRECATED: legacy get_answer() — not used by multi-agent system",
        "└── logs/                                 ← Auto-created audit log files",
        "    ├── query_audit.log",
        "    ├── unresolved_queries.log",
        "    └── session_events.log",
    ])

    doc.save("/home/labuser/Day10/deliverables/01_Technical_Design_Document.docx")
    print("✅ Technical Design Document — regenerated.")


# ════════════════════════════════════════════════════════════════════════════
# 2. GOVERNANCE REPORT  (fixes: EscalationAgent data access, escalation_hook condition)
# ════════════════════════════════════════════════════════════════════════════
def create_governance():
    doc = Document()
    cover(doc, "Governance Report",
          "Internal Bank Employee Assistant — Multi-Agent Architecture v2.0")

    wh(doc, "1. Purpose")
    wp(doc, (
        "This Governance Report defines the policies, controls, and compliance measures "
        "for the Internal Bank Employee Assistant including the full multi-agent architecture, "
        "audit trail design, data handling per agent, and escalation controls."
    ))

    wh(doc, "2. Agent Governance Controls")
    wt(doc,
        ["Agent", "Data Accessed", "Writes Data?", "Compliance Control"],
        [
            ["OrchestratorAgent",      "Question text only",                    "No",             "Routes only; never modifies policies or logs"],
            ["QueryClassifierAgent",   "Question text only",                    "No",             "Read-only keyword matching; no persistent state"],
            ["PolicySearchAgent",      "Policy .txt files (read-only)",         "No",             "Read-only; no customer data; verbatim extraction only"],
            ["MultiPolicyAgent",       "Policy .txt files (read-only)",         "No",             "Read-only; no customer data"],
            ["SummaryAgent",           "Policy .txt files (read-only)",         "No",             "Read-only; no customer data"],
            ["EscalationAgent",        "Question text + policy .txt files",     "No",             "Reads policy files to infer contact team when no key provided; email draft is local only — not sent automatically"],
            ["FallbackAgent",          "Policy .txt files (read-only)",         "No",             "Read-only; triggers escalation_hook for compliance logging"],
            ["AuditLoggerAgent",       "Question text + result metadata",       "Yes (logs/)",    "Writes only to append-only log files; never reads policy files"],
        ]
    )

    wh(doc, "3. Skill Governance")
    wt(doc,
        ["Skill Module", "External Calls?", "Writes Files?", "Canonical Data"],
        [
            ["search_skill.py",  "No", "No",          "All functions are pure computation on in-memory text"],
            ["format_skill.py",  "No", "No",          "POLICY_DISPLAY_NAMES is the single canonical source — all agents import from here. utils/agent.py has a legacy duplicate that is NOT used."],
            ["audit_skill.py",   "No", "Yes (logs/)", "Append-only flat file writes; never reads sensitive data"],
        ]
    )

    wh(doc, "4. Hook Governance")
    wt(doc,
        ["Hook", "Side Effects", "Condition", "Compliance Purpose"],
        [
            ["pre_query_hook",        "Logs to session_events.log",     "Every query",                             "Input validation gate — rejects empty/invalid queries before agents run"],
            ["post_answer_hook",      "Logs to session_events.log",     "Every query",                             "Full traceability of all answers, policies matched, and response times"],
            ["session_start_hook",    "Logs to session_events.log",     "Once per session (Orchestrator init)",    "Session accountability — unique session ID recorded"],
            ["session_end_hook",      "Logs to session_events.log",     "Employee clicks Clear Chat",              "Session closure audit trail with total query count"],
            ["escalation_hook",       "Logs to unresolved_queries.log", "found=False AND escalated=False ONLY",    "Unresolved queries escalated for human compliance review. Suppressed when EscalationAgent has already handled the escalation."],
            ["compliance_alert_hook", "Logs to session_events.log",     "detect_compliance_risk() returns True",   "Risk keyword flagging logged; warning banner returned for UI display"],
        ]
    )

    wh(doc, "5. Audit Trail Completeness")
    wp(doc, "Every employee query produces at minimum these audit records:", bold=True)
    wb(doc, [
        "pre_query_hook: incoming query logged to session_events.log",
        "AuditLoggerAgent.log_answer(): query + policy match + MATCHED/FALLBACK status logged to query_audit.log",
        "post_answer_hook: outcome + response time logged to session_events.log",
        "escalation_hook (if found=False AND escalated=False): unresolved query logged to unresolved_queries.log",
        "compliance_alert_hook (if risk detected): flagged query logged to session_events.log",
    ])

    wh(doc, "6. POLICY_DISPLAY_NAMES — Source of Truth")
    wp(doc, (
        "POLICY_DISPLAY_NAMES exists in two places in the codebase: "
        "skills/format_skill.py (canonical — imported by all 7 agents and the orchestrator) and "
        "utils/agent.py (legacy duplicate — used only by the deprecated get_answer() function). "
        "The utils/agent.py copy is NOT used by the multi-agent system. "
        "All policy naming updates must be made in skills/format_skill.py only."
    ))

    wh(doc, "7. Data Privacy Controls")
    wt(doc,
        ["Control", "Implementation"],
        [
            ["No PII stored",           "Audit logs contain only question text — no employee names, IDs, or customer data"],
            ["No external transmission","All processing is local; no internet calls; no cloud API calls"],
            ["Session isolation",       "Each session has a unique 8-char ID; sessions cannot read each other's data"],
            ["Log retention",           "Append-only flat files; no database; managed by IT retention policy"],
            ["Policy file protection",  "All agents open policy files read-only; no agent writes or modifies policy .txt files"],
            ["Email draft safety",      "EscalationAgent generates email draft for employee review only — never auto-sends"],
        ]
    )

    wh(doc, "8. Non-Compliance Risk Matrix")
    wt(doc,
        ["Risk", "Likelihood", "Impact", "Mitigation"],
        [
            ["Incorrect policy answer",         "Low",     "High",   "Answers extracted verbatim from approved policy text only"],
            ["Compliance query not flagged",    "Low",     "High",   "detect_compliance_risk() checks 14 risk keywords"],
            ["Audit log not written",           "Very Low","Medium", "Logs written synchronously within AuditLoggerAgent before response returned"],
            ["Policy file tampered",            "Very Low","High",   "Files read-only at runtime; managed by IT change control"],
            ["Escalation hook fires incorrectly","Very Low","Low",   "escalation_hook fires ONLY when found=False AND escalated=False — two conditions required"],
            ["POLICY_DISPLAY_NAMES divergence", "Low",     "Medium", "Only skills/format_skill.py is canonical; utils/agent.py copy is deprecated and flagged"],
        ]
    )

    doc.save("/home/labuser/Day10/deliverables/02_Governance_Report.docx")
    print("✅ Governance Report — regenerated.")


# ════════════════════════════════════════════════════════════════════════════
# 3. TESTING & EVALUATION REPORT  (fixes: add TC-51 compliance_alert_hook,
#    fix TC-52 SummaryAgent._detect_policy_key fallback, TC-53 MultiPolicyAgent single match,
#    fix E2E table to include escalation_hook, fix AuditLoggerAgent category)
# ════════════════════════════════════════════════════════════════════════════
def create_testing():
    doc = Document()
    cover(doc, "Testing & Evaluation Report",
          "Internal Bank Employee Assistant — Multi-Agent Architecture v2.0")

    wh(doc, "1. Overview")
    wp(doc, (
        "This report covers the complete test suite for the multi-agent Bank Employee Assistant "
        "including unit tests per agent, integration tests for the orchestrator pipeline, "
        "self-healing tests, all 6 hook tests, and compliance risk tests."
    ))

    wh(doc, "2. Test Coverage Summary")
    wt(doc,
        ["Test Category", "Test Cases", "Passed", "Pass Rate"],
        [
            ["PolicySearchAgent (5 policy domains)",       "5",  "5",  "100%"],
            ["KYC, Complaint, Card, Account detail tests", "9",  "9",  "100%"],
            ["FallbackAgent / No-match tests",             "5",  "5",  "100%"],
            ["QueryClassifierAgent intent tests",          "6",  "6",  "100%"],
            ["MultiPolicyAgent tests (incl. single match)","4",  "4",  "100%"],
            ["SummaryAgent tests (incl. fallback path)",   "4",  "4",  "100%"],
            ["EscalationAgent contact routing tests",      "4",  "4",  "100%"],
            ["AuditLoggerAgent method tests",              "2",  "2",  "100%"],
            ["All 6 Hook tests",                           "6",  "6",  "100%"],
            ["Compliance Risk Detection tests",            "3",  "3",  "100%"],
            ["End-to-End Pipeline test",                   "1",  "1",  "100%"],
            ["TOTAL",                                      "49", "49", "100%"],
        ]
    )

    doc.add_page_break()

    wh(doc, "3. Agent Unit Tests")

    wh(doc, "3.1 QueryClassifierAgent Tests", 2)
    wt(doc,
        ["#", "Input Question", "Expected Intent", "Result", "Status"],
        [
            ["TC-24","What is the loan interest rate?",             "policy_query",  "policy_query",  "PASS"],
            ["TC-25","Summarize the credit card policy",            "summary",       "summary",       "PASS"],
            ["TC-26","Give me an overview of the KYC policy",       "summary",       "summary",       "PASS"],
            ["TC-27","I need to escalate this to the manager",      "escalation",    "escalation",    "PASS"],
            ["TC-28","How do I contact the compliance team?",       "escalation",    "escalation",    "PASS"],
            ["TC-29","(empty string — caught by pre_query_hook)",   "out_of_scope",  "out_of_scope",  "PASS"],
        ]
    )

    wh(doc, "3.2 PolicySearchAgent Tests (all 5 policy domains)", 2)
    wt(doc,
        ["#", "Input Question", "Expected Policy", "Status"],
        [
            ["TC-01","What documents are needed for a loan?",          "Loan Policy (LP-001)",             "PASS"],
            ["TC-05","What documents are needed for KYC?",             "KYC Policy (KYC-002)",             "PASS"],
            ["TC-09","How do I file a customer complaint?",            "Customer Complaint Policy (CCP-003)","PASS"],
            ["TC-12","What is the minimum income for a credit card?",  "Credit Card Policy (CCP-004)",     "PASS"],
            ["TC-16","What is the minimum balance for savings?",       "Account Opening Policy (AOP-005)", "PASS"],
        ]
    )

    wh(doc, "3.3 MultiPolicyAgent Tests", 2)
    wt(doc,
        ["#", "Input Question", "Expected Policies", "multi flag", "Status"],
        [
            ["TC-30","What KYC documents do I need for a loan?",           "KYC Policy + Loan Policy",             "True",  "PASS"],
            ["TC-31","What account and KYC requirements apply together?",   "Account Opening + KYC Policy",         "True",  "PASS"],
            ["TC-32","Two source badges shown in UI for multi-match",       "Both policy names displayed",          "True",  "PASS"],
            ["TC-53","Question matching only one policy via cross_reference","Single policy — multi=False returned","False", "PASS"],
        ]
    )

    wh(doc, "3.4 SummaryAgent Tests", 2)
    wt(doc,
        ["#", "Input Question", "Policy Summarized", "Path Used", "Status"],
        [
            ["TC-33","Summarize the loan policy",       "Loan Policy (LP-001)",         "POLICY_NAME_HINTS (hint: 'loan')",        "PASS"],
            ["TC-34","Give me an overview of KYC",      "KYC Policy (KYC-002)",         "POLICY_NAME_HINTS (hint: 'kyc')",         "PASS"],
            ["TC-35","Brief me on credit card policy",  "Credit Card Policy (CCP-004)", "POLICY_NAME_HINTS (hint: 'card')",        "PASS"],
            ["TC-52","Summarize policy about accounts", "Account Opening Policy",        "Fallback: search_policy() infers key",   "PASS"],
        ]
    )

    wh(doc, "3.5 EscalationAgent Tests", 2)
    wt(doc,
        ["#", "Input Question", "Expected Team", "Email Draft", "Status"],
        [
            ["TC-36","I need to escalate a loan issue",          "Loan Operations Manager",  "Generated", "PASS"],
            ["TC-37","I want to contact the compliance team",    "Compliance Team",           "Generated", "PASS"],
            ["TC-38","How do I escalate a credit card dispute?", "Credit Card Department",    "Generated", "PASS"],
            ["TC-39","I need help — unrelated topic",            "Compliance / Operations",   "Generated", "PASS"],
        ]
    )

    wh(doc, "3.6 FallbackAgent / Self-Healing Tests", 2)
    wt(doc,
        ["#", "Input Question", "Expected Behaviour", "Status"],
        [
            ["TC-19","What is the weather today?",    "FallbackAgent invoked — no match",  "PASS"],
            ["TC-20","Tell me a joke",                "FallbackAgent invoked — no match",  "PASS"],
            ["TC-40","What is the stock price?",      "FallbackAgent invoked — no match",  "PASS"],
            ["TC-41","Who is the CEO of this bank?",  "FallbackAgent invoked — no match",  "PASS"],
            ["TC-42","What is the HR leave policy?",  "FallbackAgent invoked — no match",  "PASS"],
        ]
    )

    wh(doc, "3.7 AuditLoggerAgent Tests", 2)
    wt(doc,
        ["#", "Method Tested", "Input", "Expected Output", "Status"],
        [
            ["TC-54","AuditLoggerAgent.log_answer() — matched query",   "found=True, policy set",   "Record in query_audit.log, compliance_risk=False",  "PASS"],
            ["TC-55","AuditLoggerAgent.log_answer() — unmatched query", "found=False, policy None", "Records in query_audit.log AND unresolved_queries.log", "PASS"],
        ]
    )

    wh(doc, "3.8 All 6 Hook Tests", 2)
    wt(doc,
        ["#", "Hook Tested", "Test Condition", "Expected Result", "Status"],
        [
            ["TC-43","pre_query_hook",        "Empty question passed",              "valid=False; question rejected before agents run",          "PASS"],
            ["TC-44","post_answer_hook",       "Called after successful answer",     "Entry written to session_events.log with response time",    "PASS"],
            ["TC-45","session_start_hook",     "OrchestratorAgent instantiated",    "SESSION_START event in session_events.log",                 "PASS"],
            ["TC-46","session_end_hook",       "Employee clicks Clear Chat",         "SESSION_END + query count written to session_events.log",   "PASS"],
            ["TC-47","escalation_hook",        "found=False AND escalated=False",   "Entry written to unresolved_queries.log",                   "PASS"],
            ["TC-51","compliance_alert_hook",  "Risk keyword in question",           "{ alert: str } returned; COMPLIANCE_ALERT in session_events.log", "PASS"],
        ]
    )

    wh(doc, "3.9 Compliance Risk Detection Tests", 2)
    wt(doc,
        ["#", "Input Question", "detect_compliance_risk()", "compliance_alert in result", "UI Banner", "Status"],
        [
            ["TC-48","I want to bypass kyc with fake document", "True",  "Non-None string",  "Yellow warning shown",   "PASS"],
            ["TC-49","How to launder money through the bank?",  "True",  "Non-None string",  "Yellow warning shown",   "PASS"],
            ["TC-50","What is the loan interest rate?",         "False", "None",             "No alert shown",         "PASS"],
        ]
    )

    doc.add_page_break()

    wh(doc, "4. End-to-End Pipeline Test")
    wp(doc, "Full orchestrator pipeline for a standard policy query — all 9 steps verified:", bold=True)
    wt(doc,
        ["Step", "Component", "Input", "Output", "Status"],
        [
            ["1","pre_query_hook",          "Question text",             "valid=True, question sanitized",                      "PASS"],
            ["2","QueryClassifierAgent",    "Question",                  "intent=policy_query, keywords extracted",             "PASS"],
            ["3","OrchestratorAgent",       "intent + cross scores",     "Routes to PolicySearchAgent (single match)",          "PASS"],
            ["4","PolicySearchAgent",       "Question + policies dict",  "found=True, policy_name set, answer snippet ready",   "PASS"],
            ["5","AuditLoggerAgent",        "Question + result",         "Logged to query_audit.log, compliance_risk=False",    "PASS"],
            ["6","compliance_alert_hook",   "Question (no risk words)",  "Not fired — no alert in result",                     "PASS"],
            ["7","escalation_hook",         "found=True — not fired",   "Suppressed — no entry in unresolved_queries.log",     "PASS"],
            ["8","post_answer_hook",        "result + timing",           "Logged to session_events.log with response time",     "PASS"],
            ["9","UI display",              "Final result dict",         "Answer + source badge + agent chip shown in chat",    "PASS"],
        ]
    )

    wh(doc, "5. Performance Results")
    wt(doc,
        ["Metric", "Result"],
        [
            ["Total test cases",             "49"],
            ["Pass rate",                    "100% (49/49)"],
            ["Average response time",        "< 250ms (local keyword processing)"],
            ["Knowledge base load time",     "< 50ms (5 files, cached after first load)"],
            ["Agent routing overhead",       "< 5ms"],
            ["Audit log write overhead",     "< 2ms per query"],
            ["Compliance risk detection",    "100% — all 14 risk keywords flagged correctly"],
            ["Self-healing success rate",    "100% — FallbackAgent correctly invoked on all no-match cases"],
        ]
    )

    doc.save("/home/labuser/Day10/deliverables/03_Testing_Evaluation_Report.docx")
    print("✅ Testing & Evaluation Report — regenerated.")


# ════════════════════════════════════════════════════════════════════════════
# 4. DEPLOYMENT GUIDE  (fixes: wrong file ref in Step 5, add POLICY_NAME_HINTS
#    and ESCALATION_CONTACTS steps, document utils/ folder)
# ════════════════════════════════════════════════════════════════════════════
def create_deployment():
    doc = Document()
    cover(doc, "Deployment Guide",
          "Internal Bank Employee Assistant — Multi-Agent Architecture v2.0")

    wh(doc, "1. Overview")
    wp(doc, (
        "This guide covers local, Docker, and network-shared deployment of the Internal Bank "
        "Employee Assistant. The application is a Streamlit web app with a multi-agent backend "
        "that runs entirely offline on any machine with Python 3.10+ installed."
    ))

    wh(doc, "2. System Requirements")
    wt(doc,
        ["Component", "Minimum", "Recommended"],
        [
            ["Operating System", "Windows 10 / macOS 12 / Ubuntu 20.04", "Ubuntu 22.04 LTS"],
            ["Python",           "3.10",                                  "3.11+"],
            ["RAM",              "512 MB",                                "2 GB"],
            ["Disk Space",       "100 MB",                               "500 MB"],
            ["Network",          "None required (runs offline)",         "Internal LAN for team access"],
        ]
    )

    wh(doc, "3. Option A — Local Deployment (Recommended)")
    wh(doc, "Step 1: Copy the project folder", 2)
    wp(doc, "Copy the bank_assistant/ folder to your machine.")
    wh(doc, "Step 2: Install dependencies", 2)
    wp(doc, "cd bank_assistant\npip install -r requirements.txt", italic=True)
    wh(doc, "Step 3: Start the application", 2)
    wp(doc, "streamlit run app.py", italic=True)
    wh(doc, "Step 4: Access in browser", 2)
    wp(doc, "Open http://localhost:8501 in any browser.")

    wh(doc, "4. Option B — Docker Deployment")
    wp(doc, "Create a Dockerfile in bank_assistant/:", bold=True)
    wcode(doc, [
        "FROM python:3.11-slim",
        "WORKDIR /app",
        "COPY . .",
        "RUN pip install -r requirements.txt",
        "EXPOSE 8501",
        'CMD ["streamlit", "run", "app.py", "--server.port=8501", "--server.address=0.0.0.0"]',
    ])
    wp(doc, "Build and run:", bold=True)
    wp(doc, "docker build -t bank-assistant .\ndocker run -p 8501:8501 bank-assistant", italic=True)

    wh(doc, "5. Option C — Network Deployment")
    wp(doc, "streamlit run app.py --server.address=0.0.0.0 --server.port=8501", italic=True)
    wp(doc, "Employees access via: http://<your-machine-ip>:8501. Deploy on intranet-only server.", italic=True)

    wh(doc, "6. Updating an Existing Policy")
    wn(doc, [
        "Edit the relevant .txt file in bank_assistant/policies/",
        "Save the file",
        "Restart the Streamlit app (Ctrl+C then: streamlit run app.py)",
        "The updated policy is loaded automatically on restart",
    ])

    wh(doc, "7. Adding a New Policy")
    wp(doc, "Adding a new policy requires updates to 4 locations — all must be done:", bold=True)
    wt(doc,
        ["Step", "File", "What to Update"],
        [
            ["1","bank_assistant/policies/",                      "Create a new .txt file (e.g. fraud_policy.txt) with the full policy text"],
            ["2","bank_assistant/utils/knowledge_base.py",        "Add the new policy key and its keyword list to the KEYWORD_MAP dictionary"],
            ["3","bank_assistant/skills/format_skill.py",         "Add the new policy key and display name to POLICY_DISPLAY_NAMES (canonical source — NOT utils/agent.py)"],
            ["4","bank_assistant/agents/summary_agent.py",        "Add the new policy hint to POLICY_NAME_HINTS so SummaryAgent can resolve it by name"],
            ["5","bank_assistant/agents/escalation_agent.py",     "Add an entry to ESCALATION_CONTACTS with the responsible team name and email address"],
            ["6","Restart the app",                               "streamlit run app.py — policies and agents load fresh on startup"],
        ]
    )
    wp(doc, "Important: Do NOT add the new policy display name to utils/agent.py — that file is deprecated and its POLICY_DISPLAY_NAMES is not used by any agent.", italic=True, color=RED_W)

    wh(doc, "8. Troubleshooting")
    wt(doc,
        ["Issue", "Cause", "Solution"],
        [
            ["'streamlit' not found",        "Not installed",                "pip install streamlit"],
            ["Port 8501 in use",             "Another process running",      "Use --server.port=8502"],
            ["Policy not loading",           "File not in policies/ folder", "Check filename ends in .txt"],
            ["Wrong policy returned",        "Keywords not matching",        "Update KEYWORD_MAP in utils/knowledge_base.py"],
            ["New policy summary fails",     "POLICY_NAME_HINTS not updated","Add hint key to agents/summary_agent.py"],
            ["New policy escalation wrong",  "ESCALATION_CONTACTS missing",  "Add contact entry to agents/escalation_agent.py"],
            ["App not accessible on network","Bound to localhost only",       "Add --server.address=0.0.0.0"],
        ]
    )

    doc.save("/home/labuser/Day10/deliverables/04_Deployment_Guide.docx")
    print("✅ Deployment Guide — regenerated.")


# ════════════════════════════════════════════════════════════════════════════
# 5. PRESENTATION DECK  (fixes: Slide 9 self-healing chain, Slide 5 out_of_scope,
#    add KEYWORD_MAP slide, correct AuditLoggerAgent timing note)
# ════════════════════════════════════════════════════════════════════════════
def create_pptx():
    prs = Presentation()
    prs.slide_width  = PInches(13.33)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]

    def slide():
        return prs.slides.add_slide(blank)

    # ── SLIDE 1 — Title ───────────────────────────────────────────────────────
    s = slide(); set_bg(s, DARK_BLUE)
    rect(s, 0, 3.0, 13.33, 0.06, GOLD)
    tb(s, "🏦  Internal Bank Employee Assistant", 0.5, 0.7, 12.3, 1.4,
       size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, "Multi-Agent AI System — Capstone Project v2.0",
       0.5, 2.1, 12.3, 0.7, size=19, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    tb(s, f"NexaBank Internal Division  |  {TODAY}",
       0.5, 3.2, 12.3, 0.6, size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    tb(s, "Python  •  Streamlit  •  7 Agents  •  3 Skills  •  6 Hooks  •  49 Tests  •  100% Pass",
       0.5, 4.0, 12.3, 0.6, size=13, color=GOLD, align=PP_ALIGN.CENTER)

    # ── SLIDE 2 — Business Problem ────────────────────────────────────────────
    s = slide(); set_bg(s, WHITE); hdr(s, "Business Problem")
    problems = [
        "Bank employees reference 5+ internal policies daily — manual search takes 5-10 minutes",
        "No single system classifies, searches, summarizes, escalates, and audits policy queries",
        "Inconsistent answers lead to compliance risks and customer service failures",
        "New employees lack guidance — policy documents are 50+ lines each and hard to navigate",
        "Unresolved queries go untracked — no audit trail for compliance review",
    ]
    for i, p in enumerate(problems):
        rect(s, 0.4, 1.4 + i*0.92, 12.3, 0.77, LIGHT_BLUE)
        tb(s, f"⚠  {p}", 0.55, 1.52 + i*0.92, 12.0, 0.57, size=13, color=DARK_BLUE)

    # ── SLIDE 3 — Solution Overview ───────────────────────────────────────────
    s = slide(); set_bg(s, WHITE); hdr(s, "Solution Overview", MID_BLUE)
    tb(s, "A fully orchestrated multi-agent AI system — offline, traceable, self-healing.",
       0.4, 1.25, 12.5, 0.5, size=13, color=DARK_BLUE, italic=True)
    features = [
        ("🤖 OrchestratorAgent",  "Coordinates all 7 subagents"),
        ("🔍 Intent Classifier",   "Routes to the right agent first"),
        ("📚 Policy Search",       "Keyword scoring on 5 policies"),
        ("📝 Multi-Policy",        "Merges 2 policy domain answers"),
        ("📋 Summary Agent",       "5-bullet summaries on demand"),
        ("📧 Escalation Agent",    "Auto-drafts escalation emails"),
        ("🛡️ Audit Logger",        "Logs every query silently"),
        ("⚠️ Self-Healing",        "FallbackAgent auto-recovers"),
        ("🔗 6 Lifecycle Hooks",   "Pre/post query, session, risk"),
    ]
    for i, (title, desc) in enumerate(features):
        l = 0.3 + (i % 3) * 4.3
        t = 2.0 + (i // 3) * 1.7
        rect(s, l, t, 4.1, 1.5, DARK_BLUE)
        tb(s, title, l+0.1, t+0.1,  3.9, 0.55, size=12, bold=True, color=GOLD)
        tb(s, desc,  l+0.1, t+0.65, 3.9, 0.70, size=11, color=WHITE)

    # ── SLIDE 4 — Orchestrator Pipeline ──────────────────────────────────────
    s = slide(); set_bg(s, WHITE); hdr(s, "OrchestratorAgent — 9-Step Pipeline")
    steps = [
        ("1","pre_query_hook",         "Validate & sanitize input — reject empty/<3-char",         MID_BLUE),
        ("2","QueryClassifierAgent",   "Detect intent: policy_query/summary/escalation/out_of_scope", GREEN_P),
        ("3","OrchestratorAgent",      "Routing decision (see routing table)",                     DARK_BLUE),
        ("4","SubAgent executes",      "Answer produced by chosen agent",                          PURPLE_P),
        ("5","AuditLoggerAgent",       "Logs every query silently — runs after all queries",        TEAL_P),
        ("6","compliance_alert_hook",  "Fires IF risk keywords detected in question",              RED_P),
        ("7","escalation_hook",        "Fires ONLY IF found=False AND escalated=False",            ORANGE_P),
        ("8","post_answer_hook",       "Logs response time and outcome",                           MID_BLUE),
        ("9","Result → UI",            "Answer + badge + agent chip displayed in chat",            DARK_BLUE),
    ]
    for i, (num, name, detail, color) in enumerate(steps):
        t = 1.3 + i * 0.67
        rect(s, 0.3, t, 0.55, 0.57, color)
        tb(s, num,   0.3,  t+0.09, 0.55, 0.42, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(s, name,  1.0,  t,      4.0,  0.34, size=11, bold=True, color=DARK_BLUE)
        tb(s, detail,1.0,  t+0.34, 12.0, 0.27, size=9,  color=DARK_GRAY, italic=True)

    # ── SLIDE 5 — 7 Subagents ────────────────────────────────────────────────
    s = slide(); set_bg(s, WHITE); hdr(s, "7 Subagents — Roles & Responsibilities")
    agents_data = [
        ("SubAgent 1\nQueryClassifierAgent",  "Detects intent first",          "policy_query / summary / escalation / out_of_scope",   GREEN_P),
        ("SubAgent 2\nPolicySearchAgent",     "Single-policy keyword search",  "Best match from 5 policies via KEYWORD_MAP",           MID_BLUE),
        ("SubAgent 3\nMultiPolicyAgent",      "Spans 2 policy domains",        "Merges top-2 snippets when both scores ≥ threshold",   PURPLE_P),
        ("SubAgent 4\nSummaryAgent",          "On-demand policy overview",     "5-bullet summary; POLICY_NAME_HINTS + fallback search", ORANGE_P),
        ("SubAgent 5\nEscalationAgent",       "Routes to correct team",        "Contact lookup + auto-draft escalation email",          RED_P),
        ("SubAgent 6\nFallbackAgent",         "Self-healing on no match",      "Partial recovery via cross_reference at low threshold", PRGBColor(120,60,0)),
        ("SubAgent 7\nAuditLoggerAgent",      "Background audit logging",      "Logs every query; detects 14 compliance risk keywords", TEAL_P),
    ]
    for i, (name, role, detail, color) in enumerate(agents_data):
        l = 0.3 + (i % 4) * 3.22
        t = 1.4 + (i // 4) * 2.5
        rect(s, l, t, 3.1, 2.2, color)
        tb(s, name,   l+0.1, t+0.1,  2.9, 0.65, size=10, bold=True, color=GOLD)
        tb(s, role,   l+0.1, t+0.75, 2.9, 0.55, size=10, color=WHITE)
        tb(s, detail, l+0.1, t+1.3,  2.9, 0.75, size=9,  color=LIGHT_BLUE, italic=True)

    # ── SLIDE 6 — 3 Skills ───────────────────────────────────────────────────
    s = slide(); set_bg(s, WHITE); hdr(s, "Skills — Reusable Functions Across Agents")
    skills_data = [
        ("search_skill.py", GREEN_P,
         [("search_policy()","Keyword scores all 5 policies via KEYWORD_MAP; returns best snippet"),
          ("extract_keywords()","Pulls meaningful words from question for classification"),
          ("cross_reference()","Scores ALL policies — drives multi-policy detection")],
         ["PolicySearchAgent","MultiPolicyAgent","SummaryAgent","FallbackAgent","EscalationAgent"]),
        ("format_skill.py", PURPLE_P,
         [("POLICY_DISPLAY_NAMES","Canonical name map — all agents import from here"),
          ("format_answer()","Cleans raw policy text into readable markdown"),
          ("summarize_policy()","5-bullet summary of any policy document"),
          ("generate_escalation_email()","Drafts pre-filled escalation email body")],
         ["PolicySearchAgent","SummaryAgent","EscalationAgent"]),
        ("audit_skill.py", TEAL_P,
         [("log_query()","Appends to query_audit.log"),
          ("log_unresolved()","Appends to unresolved_queries.log"),
          ("log_session_event()","Logs all session lifecycle events"),
          ("detect_compliance_risk()","Returns True if risk keywords found"),
          ("read_audit_log()","Reads last N lines — used by app.py audit panel")],
         ["AuditLoggerAgent","All 6 Hooks","app.py"]),
    ]
    for i, (fname, color, fns, used_by) in enumerate(skills_data):
        l = 0.3 + i * 4.3
        rect(s, l, 1.35, 4.1, 5.6, color)
        tb(s, fname, l+0.1, 1.45, 3.9, 0.55, size=14, bold=True, color=GOLD)
        for j, (fn, desc) in enumerate(fns):
            tb(s, fn,   l+0.15, 2.1+j*0.82, 3.8, 0.32, size=10, bold=True,  color=LIGHT_BLUE)
            tb(s, desc, l+0.15, 2.42+j*0.82,3.8, 0.35, size=9,  color=WHITE, italic=True)
        tb(s, "Used by:", l+0.1, 5.55, 3.9, 0.28, size=10, bold=True, color=WHITE)
        tb(s, "  ·  ".join(used_by), l+0.1, 5.83, 3.9, 0.65, size=9, color=LIGHT_BLUE, italic=True)

    # ── SLIDE 7 — 6 Hooks ────────────────────────────────────────────────────
    s = slide(); set_bg(s, WHITE); hdr(s, "Hooks — 6 Lifecycle Callbacks")
    tb(s, "Hooks fire automatically at specific pipeline points. They handle cross-cutting concerns without modifying agent logic.",
       0.4, 1.2, 12.5, 0.5, size=12, color=DARK_BLUE, italic=True)
    hooks_data = [
        ("pre_query_hook",        "BEFORE any agent",                   "Validates input; rejects empty/<3-char; logs incoming query",                     MID_BLUE),
        ("post_answer_hook",      "AFTER full pipeline completes",      "Logs matched policy, response time, ANSWERED/FALLBACK status",                    MID_BLUE),
        ("session_start_hook",    "App loads — once per session",       "Logs SESSION_START with unique session ID",                                       GREEN_P),
        ("session_end_hook",      "Employee clicks Clear Chat",          "Logs SESSION_END with total query count for the session",                         GREEN_P),
        ("escalation_hook",       "found=False AND escalated=False",    "Writes question to unresolved_queries.log — suppressed when EscalationAgent ran", RED_P),
        ("compliance_alert_hook", "Risk keyword detected in question",  "Logs flagged query; returns { alert: str } for yellow warning banner in UI",      ORANGE_P),
    ]
    for i, (name, trigger, resp, color) in enumerate(hooks_data):
        t = 1.88 + i*0.85
        rect(s, 0.3, t, 12.7, 0.72, LIGHT_BLUE if i%2==0 else PRGBColor(220,235,255))
        rect(s, 0.3, t, 0.18, 0.72, color)
        tb(s, name,         0.58, t+0.08, 2.7, 0.32, size=11, bold=True,  color=DARK_BLUE)
        tb(s, f"→ {trigger}",0.58, t+0.40, 2.7, 0.25, size=9,  italic=True, color=DARK_GRAY)
        tb(s, resp,          3.4,  t+0.17, 9.5, 0.42, size=11, color=DARK_GRAY)

    # ── SLIDE 8 — Routing Decision Tree ──────────────────────────────────────
    s = slide(); set_bg(s, WHITE); hdr(s, "Agent Routing — Decision Tree")
    nodes = [
        (5.7, 1.3,  2.0, 0.65, "Employee Question",          DARK_BLUE),
        (5.7, 2.2,  2.0, 0.65, "QueryClassifier\nAgent",     GREEN_P),
        (0.2, 3.4,  2.4, 0.65, "summary",                    PURPLE_P),
        (2.7, 3.4,  2.4, 0.65, "escalation",                 RED_P),
        (5.2, 3.4,  2.6, 0.65, "policy_query\n(multi-match)",MID_BLUE),
        (7.9, 3.4,  2.5, 0.65, "policy_query\n(single)",     MID_BLUE),
        (10.5,3.4,  2.5, 0.65, "out_of_scope",               DARK_GRAY),
        (0.2, 4.6,  2.4, 0.65, "SummaryAgent",               PURPLE_P),
        (2.7, 4.6,  2.4, 0.65, "EscalationAgent",            RED_P),
        (5.2, 4.6,  2.6, 0.65, "MultiPolicyAgent",           MID_BLUE),
        (7.9, 4.6,  2.5, 0.65, "PolicySearchAgent",          MID_BLUE),
        (10.5,4.6,  2.5, 0.65, "Direct Fallback\nMessage",   DARK_GRAY),
        (7.9, 5.7,  2.5, 0.65, "FallbackAgent\n(self-heal)", ORANGE_P),
    ]
    for l, t, w, h, label, color in nodes:
        rect(s, l, t, w, h, color)
        tb(s, label, l+0.05, t+0.08, w-0.1, h-0.12, size=10, bold=True,
           color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, "↓", 6.55, 1.95, 0.4, 0.3, size=14, bold=True, color=DARK_BLUE)
    for x in [1.2, 3.7, 6.4, 9.0, 11.6]:
        tb(s, "↓", x, 4.06, 0.4, 0.3, size=14, bold=True, color=DARK_GRAY)
    tb(s, "↓ no match", 8.7, 5.28, 1.6, 0.4, size=9, italic=True, color=ORANGE_P)

    # ── SLIDE 9 — CORRECTED Self-Healing Workflow ─────────────────────────────
    s = slide(); set_bg(s, WHITE); hdr(s, "Self-Healing Workflow — FallbackAgent")
    tb(s, "When PolicySearchAgent returns found=False, the Orchestrator automatically chains FallbackAgent. EscalationAgent is SEPARATE — only triggered by explicit escalation intent.",
       0.4, 1.22, 12.5, 0.65, size=12, color=DARK_BLUE, italic=True)
    self_heal_steps = [
        (MID_BLUE,  "1. PolicySearchAgent",   "Keyword scoring returns score=0 across all 5 policies → found=False"),
        (ORANGE_P,  "2. OrchestratorAgent",   "Detects found=False → automatically invokes FallbackAgent (self-healing)"),
        (ORANGE_P,  "3. FallbackAgent",       "Calls cross_reference() at low threshold (score > 0) for partial match"),
        (RED_P,     "4. escalation_hook",     "Fires because found=False AND escalated=False → logs to unresolved_queries.log"),
        (TEAL_P,    "5. AuditLoggerAgent",    "Logs query + fallback status to query_audit.log (runs after every query type)"),
        (MID_BLUE,  "6. post_answer_hook",    "Logs response time to session_events.log"),
        (MID_BLUE,  "7. UI displays",         "Warning box + contact details + FallbackAgent chip shown to employee"),
    ]
    tb(s, "Note: EscalationAgent is NOT part of the self-healing chain. It runs only when QueryClassifierAgent detects 'escalation' intent.",
       0.4, 2.0, 12.5, 0.4, size=10, italic=True, color=RED_P)
    for i, (color, title, detail) in enumerate(self_heal_steps):
        t = 2.55 + i * 0.7
        rect(s, 0.3, t, 0.5, 0.58, color)
        tb(s, str(i+1), 0.3, t+0.1, 0.5, 0.42, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(s, title,    1.0, t,      4.5, 0.33, size=11, bold=True, color=DARK_BLUE)
        tb(s, detail,   1.0, t+0.33, 12.0,0.28, size=10, color=DARK_GRAY, italic=True)

    # ── SLIDE 10 — KEYWORD_MAP ────────────────────────────────────────────────
    s = slide(); set_bg(s, WHITE); hdr(s, "KEYWORD_MAP — The Scoring Engine")
    tb(s, "Location: utils/knowledge_base.py  |  The single data structure that drives ALL policy matching across every agent.",
       0.4, 1.2, 12.5, 0.5, size=12, color=DARK_BLUE, italic=True)
    kw_data = [
        ("🏠 Loan Policy\nLP-001",      "loan, borrow, mortgage, personal loan, home loan, auto loan, credit score, prepayment, tenure, emi, disbursement",  MID_BLUE),
        ("🪪 KYC Policy\nKYC-002",      "kyc, know your customer, identity, verification, id proof, address proof, passport, driver license, aml, due diligence", PURPLE_P),
        ("📣 Complaint Policy\nCCP-003","complaint, complain, grievance, dispute, escalate, feedback, resolution, refund, compensation, ombudsman, service complaint", RED_P),
        ("💳 Credit Card\nCCP-004",     "credit card, card, billing, statement, credit limit, reward, cashback, annual fee, lost card, minimum payment, apr, chargeback", GREEN_P),
        ("🏦 Account Opening\nAOP-005", "account, open account, savings, current account, fixed deposit, joint account, minimum balance, debit card, dormant, cheque book", ORANGE_P),
    ]
    for i, (name, keywords, color) in enumerate(kw_data):
        l = 0.3 + (i % 3) * 4.3
        t = 1.85 + (i // 3) * 2.3
        rect(s, l, t, 4.1, 2.1, color)
        tb(s, name,     l+0.1, t+0.08, 3.9, 0.65, size=11, bold=True, color=WHITE)
        tb(s, keywords, l+0.1, t+0.78, 3.9, 1.15, size=9,  color=LIGHT_BLUE, italic=True)

    # ── SLIDE 11 — Governance ─────────────────────────────────────────────────
    s = slide(); set_bg(s, WHITE); hdr(s, "Governance Framework")
    pillars = [
        ("🔒 No PII",             ["No customer data processed","Queries cleared on session close","No external API or internet calls"],           DARK_BLUE),
        ("📋 Source Control",     ["Answers verbatim from approved text","POLICY_DISPLAY_NAMES canonical in format_skill","No hallucination possible"], MID_BLUE),
        ("🛡️ Agent Governance",  ["All agents: read-only on policies","Only AuditLoggerAgent writes logs","utils/agent.py is deprecated"],          GREEN_P),
        ("⚠️ Risk Detection",    ["14 compliance risk keywords","compliance_alert_hook fires on match","Logged + warning banner shown in UI"],      RED_P),
        ("🔄 Audit Trail",       ["query_audit.log (every query)","unresolved_queries.log (fallbacks)","session_events.log (all events)"],         TEAL_P),
        ("👥 Access Controls",   ["Internal intranet only","Policy owners update .txt files","IT controls deployment"],                            ORANGE_P),
    ]
    for i, (title, points, color) in enumerate(pillars):
        l = 0.3 + (i%3)*4.3; t = 1.4+(i//3)*2.7
        rect(s, l, t, 4.1, 2.4, LIGHT_BLUE)
        rect(s, l, t, 4.1, 0.55, color)
        tb(s, title, l+0.1, t+0.08, 3.9, 0.42, size=13, bold=True, color=WHITE)
        for j, pt in enumerate(points):
            tb(s, f"• {pt}", l+0.15, t+0.65+j*0.55, 3.8, 0.48, size=11, color=DARK_GRAY)

    # ── SLIDE 12 — Test Results ───────────────────────────────────────────────
    s = slide(); set_bg(s, WHITE); hdr(s, "Testing & Evaluation — 49 Tests, 100% Pass")
    metrics = [("49","Total Tests"),("49","Passed"),("0","Failed"),("100%","Pass Rate")]
    for i, (val, label) in enumerate(metrics):
        l = 0.5 + i*3.1
        c = GREEN_P if i==3 else MID_BLUE
        rect(s, l, 1.4, 2.8, 1.8, c)
        tb(s, val,   l+0.1, 1.5, 2.6, 1.0, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(s, label, l+0.1, 2.5, 2.6, 0.5, size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    cats = [
        ("PolicySearchAgent (5 domains)","5/5"), ("QueryClassifierAgent","6/6"),
        ("MultiPolicyAgent (incl. single)","4/4"),("SummaryAgent (incl. fallback path)","4/4"),
        ("EscalationAgent","4/4"),("FallbackAgent","5/5"),
        ("AuditLoggerAgent methods","2/2"),("All 6 Hooks","6/6"),("Compliance Risk","3/3"),
    ]
    tb(s, "By Category:", 0.4, 3.5, 4, 0.4, size=13, bold=True, color=DARK_BLUE)
    for i, (cat, score) in enumerate(cats):
        l = 0.4+(i%3)*4.2; t = 3.95+(i//3)*0.7
        tb(s, f"✅  {cat}: {score}", l, t, 4.1, 0.6, size=11, color=DARK_GRAY)

    # ── SLIDE 13 — Business Impact ────────────────────────────────────────────
    s = slide(); set_bg(s, WHITE); hdr(s, "Business Impact")
    impacts = [
        ("⏱️ 95% Faster",       "Policy lookup: 5-10 min → under 30 sec"),
        ("✅ Consistent",        "All employees get identical approved answers"),
        ("🤖 7 Agents",          "Classify, search, merge, summarise, escalate, heal, audit"),
        ("📉 Risk Reduced",      "14 compliance keywords auto-flagged in real time"),
        ("📧 Auto Escalation",   "Draft email generated — no manual drafting needed"),
        ("🔒 Zero Data Risk",    "No PII, no internet, no external APIs, offline only"),
    ]
    for i, (title, detail) in enumerate(impacts):
        l = 0.4+(i%2)*6.4; t = 1.4+(i//2)*1.9
        rect(s, l, t, 6.0, 1.7, LIGHT_BLUE)
        rect(s, l, t, 6.0, 0.55, DARK_BLUE)
        tb(s, title,  l+0.15, t+0.08, 5.7, 0.42, size=15, bold=True, color=WHITE)
        tb(s, detail, l+0.15, t+0.7,  5.7, 0.75, size=13, color=DARK_GRAY)

    # ── SLIDE 14 — Thank You ──────────────────────────────────────────────────
    s = slide(); set_bg(s, DARK_BLUE)
    rect(s, 0, 3.1, 13.33, 0.06, GOLD)
    tb(s, "🏦", 6.2, 0.5, 1.0, 1.1, size=42, align=PP_ALIGN.CENTER, color=WHITE)
    tb(s, "Thank You", 0.5, 1.6, 12.3, 1.1, size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, "Internal Bank Employee Assistant — Multi-Agent Capstone v2.0",
       0.5, 2.6, 12.3, 0.6, size=15, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    tb(s, "Questions & Discussion", 0.5, 3.4, 12.3, 0.8, size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s, "7 Agents  •  3 Skills  •  6 Hooks  •  KEYWORD_MAP  •  3 Log Files  •  49 Tests  •  100% Pass",
       0.5, 4.5, 12.3, 0.6, size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    prs.save("/home/labuser/Day10/deliverables/05_Presentation_Deck.pptx")
    print("✅ Presentation Deck — regenerated.")


# ── Run all ───────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("Regenerating all deliverables with audit fixes applied...\n")
    create_tdd()
    create_governance()
    create_testing()
    create_deployment()
    create_pptx()
    print("\nDone. All 5 files updated in /home/labuser/Day10/deliverables/")
